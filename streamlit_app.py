import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import kaggle
import os
from datetime import datetime
import json
import requests
from dotenv import load_dotenv
import base64
import io
from PIL import Image
import tempfile
import warnings
import random

# Suppress all warnings for cleaner output
warnings.filterwarnings('ignore')

# Optional imports that may not be available
try:
    import numpy as np
    import scipy.stats  # Changed from 'from scipy import stats' to avoid conflicts
    HAS_SCIPY = True
except ImportError:
    HAS_SCIPY = False
    st.sidebar.info("💡 Install scipy for statistical trend lines: pip install scipy")

# Load environment variables
load_dotenv()

# Page configuration
st.set_page_config(
    page_title="KPI Insights Dashboard",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Friendly progress messages
friendly_progress_messages = [
    "🎨 Creating your visual story...",
    "📊 Bringing your data to life...",
    "✨ Making sense of your numbers...",
    "🔍 Finding the interesting patterns...",
    "🎯 Highlighting what matters most...",
    "📈 Building your insights dashboard..."
]

def create_friendly_chart_title(opportunity_title, df, opportunity):
    """Create more engaging chart titles"""
    base_title = opportunity_title.replace('_', ' ').title()
    
    # Make titles more conversational
    title_variations = {
        'Distribution': [
            f"How Your {base_title.replace('Distribution', '').strip()} is Spread Out",
            f"The Story Behind Your {base_title.replace('Distribution', '').strip()} Numbers",
            f"Breaking Down Your {base_title.replace('Distribution', '').strip()}"
        ],
        'Over Time': [
            f"Your {base_title.replace('Over Time', '').strip()} Journey",
            f"How {base_title.replace('Over Time', '').strip()} Changed Over Time",
            f"The {base_title.replace('Over Time', '').strip()} Timeline"
        ],
        'Trend': [
            f"Your {base_title.replace('Trend', '').strip()} Journey Over Time",
            f"How {base_title.replace('Trend', '').strip()} Evolved",
            f"The {base_title.replace('Trend', '').strip()} Story"
        ]
    }
    
    # Handle 'by' pattern more safely
    if 'by' in base_title.lower():
        try:
            parts = base_title.split(' by ')
            if len(parts) >= 2:
                before_by = parts[0].strip()
                after_by = parts[1].strip()
                by_variations = [
                    f"Which {after_by} Performs Best in {before_by}?",
                    f"Comparing {before_by} Across {after_by}",
                    f"Your {after_by} Performance Breakdown"
                ]
                return random.choice(by_variations)
        except:
            pass  # Fall back to original title
    
    # Find matching pattern and return random variation
    for pattern, variations in title_variations.items():
        if pattern.lower() in base_title.lower():
            return random.choice(variations)
    
    return base_title

def create_friendly_annotations():
    """Get friendly annotation text"""
    annotation_styles = {
        'high_point': ["📈 Peak performance!", "🚀 This is your best!", "⭐ Top result", "🏆 Your winner!"],
        'low_point': ["📉 Room to grow", "🎯 Focus area", "💡 Opportunity here", "⚠️ Needs attention"],
        'trend_up': ["📈 Going up!", "🚀 Great momentum", "⬆️ Rising trend", "💪 Growing strong"],
        'trend_down': ["📉 Declining", "⚠️ Needs attention", "📊 Dropping off", "🔍 Worth investigating"]
    }
    
    return annotation_styles

# Move this function to the top to avoid import issues
def identify_chart_opportunities(df, use_full_dataset=True, sample_size=2000):
    """Identify the best chart types based on data patterns"""
    try:
        # Basic validation
        if df is None or df.empty:
            return []
        
        if len(df.columns) < 1:
            return []
        
        numeric_cols = df.select_dtypes(include=['number']).columns.tolist()
        categorical_cols = df.select_dtypes(include=['object', 'category']).columns.tolist()
        datetime_cols = []
        
        # Try to identify datetime columns that are stored as strings (with NO warnings)
        for col in categorical_cols.copy():
            if df[col].dtype == 'object':
                try:
                    # Sample just 5 rows and suppress all warnings
                    sample = df[col].dropna().head(5)
                    if len(sample) > 0:
                        # Try parsing with specific formats first to avoid dateutil warnings
                        with warnings.catch_warnings():
                            warnings.simplefilter("ignore")
                            pd.to_datetime(sample, errors='raise')
                        datetime_cols.append(col)
                        categorical_cols.remove(col)
                except:
                    # Not a datetime column, continue silently
                    pass
        
        opportunities = []
        
        # Handle dataset sampling based on user preference
        try:
            if use_full_dataset or len(df) <= sample_size:
                df_sample = df
            else:
                df_sample = df.sample(n=min(sample_size, len(df)), random_state=42)
        except Exception as e:
            df_sample = df  # Use full dataset if sampling fails
        
        # Time series opportunities (create more if available)
        if len(datetime_cols) > 0 and len(numeric_cols) > 0:
            # Create time series for first 2 numeric columns with datetime
            for i, num_col in enumerate(numeric_cols[:2]):
                for j, date_col in enumerate(datetime_cols[:1]):  # Use first datetime column
                    try:
                        friendly_title = create_friendly_chart_title(f'{num_col.replace("_", " ").title()} Over Time', df_sample, None)
                        opportunities.append({
                            'type': 'time_series',
                            'title': friendly_title,
                            'x': date_col,
                            'y': num_col,
                            'priority': 'high',
                            'data': df_sample
                        })
                        if len(opportunities) >= 8:  # Limit total opportunities
                            break
                    except Exception as e:
                        continue  # Skip this opportunity if there's an error
                if len(opportunities) >= 8:
                    break
        
        # KPI metric opportunities - create more distributions
        if len(numeric_cols) > 0:
            count = 0
            for col in numeric_cols[:6]:  # Check first 6 columns instead of 5
                try:
                    # Look for KPI-like column names
                    kpi_keywords = ['revenue', 'sales', 'profit', 'cost', 'price', 'amount', 'value', 
                                   'count', 'total', 'avg', 'mean', 'sum', 'rate', 'percent', 'ctr', 'clicks', 'impressions']
                    
                    if any(keyword in col.lower() for keyword in kpi_keywords) and count < 4:  # Increased to 4
                        friendly_title = create_friendly_chart_title(f'{col.replace("_", " ").title()} Distribution', df_sample, None)
                        opportunities.append({
                            'type': 'kpi_metric',
                            'title': friendly_title,
                            'column': col,
                            'priority': 'high',
                            'data': df_sample
                        })
                        count += 1
                        if len(opportunities) >= 8:
                            break
                except Exception as e:
                    continue  # Skip this opportunity if there's an error
                if len(opportunities) >= 8:
                    break
        
        # Category performance - create multiple category analyses
        if len(categorical_cols) > 0 and len(numeric_cols) > 0:
            chart_count = 0
            for cat_col in categorical_cols[:3]:  # Check first 3 categorical columns
                try:
                    unique_count = df[cat_col].nunique()
                    if 3 <= unique_count <= 15 and chart_count < 3:  # Create up to 3 category charts
                        for num_col in numeric_cols[:2]:  # Use first 2 numeric columns
                            if chart_count < 3 and len(opportunities) < 8:
                                try:
                                    friendly_title = create_friendly_chart_title(f'{num_col.replace("_", " ").title()} by {cat_col.replace("_", " ").title()}', df_sample, None)
                                    opportunities.append({
                                        'type': 'category_performance',
                                        'title': friendly_title,
                                        'x': cat_col,
                                        'y': num_col,
                                        'priority': 'medium',
                                        'data': df_sample
                                    })
                                    chart_count += 1
                                except Exception as e:
                                    continue  # Skip this opportunity if there's an error
                            if len(opportunities) >= 8:
                                break
                        if len(opportunities) >= 8:
                            break
                except Exception as e:
                    continue  # Skip this category if there's an error
                if len(opportunities) >= 8:
                    break
        
        # Correlation analysis
        if 3 <= len(numeric_cols) <= 8 and len(opportunities) < 8:
            try:
                opportunities.append({
                    'type': 'correlation',
                    'title': 'How Your Key Metrics Connect to Each Other',
                    'columns': numeric_cols[:6],  # Up to 6 columns for correlation
                    'priority': 'medium',
                    'data': df_sample
                })
            except Exception as e:
                pass  # Skip correlation if error
        
        # Add geographic analysis if detected
        if len(opportunities) < 8:
            try:
                geo_keywords = ['country', 'state', 'city', 'region', 'location', 'address', 'area']
                geo_cols = [col for col in categorical_cols if any(keyword in col.lower() for keyword in geo_keywords)]
                
                if len(geo_cols) > 0 and len(numeric_cols) > 0:
                    for geo_col in geo_cols[:2]:  # Up to 2 geographic columns
                        for num_col in numeric_cols[:2]:  # With first 2 numeric columns
                            if len(opportunities) >= 8:
                                break
                            try:
                                if df[geo_col].nunique() <= 20:  # Reasonable number of locations
                                    friendly_title = create_friendly_chart_title(f'{num_col.replace("_", " ").title()} by {geo_col.replace("_", " ").title()}', df_sample, None)
                                    opportunities.append({
                                        'type': 'geographic',
                                        'title': friendly_title,
                                        'x': geo_col,
                                        'y': num_col,
                                        'priority': 'medium',
                                        'data': df_sample
                                    })
                            except Exception as e:
                                continue  # Skip this opportunity if there's an error
                        if len(opportunities) >= 8:
                            break
            except Exception as e:
                pass  # Skip geographic analysis if error
        
        # If no opportunities found, create basic ones
        if len(opportunities) == 0:
            try:
                # At least try to create a simple histogram if we have numeric data
                if len(numeric_cols) > 0:
                    col = numeric_cols[0]
                    opportunities.append({
                        'type': 'kpi_metric',
                        'title': f'Distribution of {col.replace("_", " ").title()}',
                        'column': col,
                        'priority': 'high',
                        'data': df_sample
                    })
                
                # Or a simple category count if we have categorical data
                elif len(categorical_cols) > 0:
                    col = categorical_cols[0]
                    if df[col].nunique() <= 20:
                        opportunities.append({
                            'type': 'category_count',
                            'title': f'Count of {col.replace("_", " ").title()}',
                            'column': col,
                            'priority': 'medium',
                            'data': df_sample
                        })
            except Exception as e:
                pass  # Even basic opportunities failed
        
        # Limit to maximum 8 charts
        return opportunities[:8]
        
    except Exception as e:
        st.error(f"Error in identify_chart_opportunities: {str(e)}")
        return []

class SmartChartGenerator:
    """Generate business-relevant charts based on data patterns"""
    
    @staticmethod
    def create_chart(df, opportunity):
        """Create fast, professional charts with key highlights"""
        try:
            # Use the sampled data from opportunity for faster rendering
            data_to_use = opportunity.get('data', df)
            annotations = create_friendly_annotations()
            
            if opportunity['type'] == 'time_series':
                # Fast time series chart with highlights
                df_copy = data_to_use.copy()
                try:
                    with warnings.catch_warnings():
                        warnings.simplefilter("ignore")
                        df_copy[opportunity['x']] = pd.to_datetime(df_copy[opportunity['x']], errors='coerce')
                    df_copy = df_copy.dropna(subset=[opportunity['x']]).sort_values(opportunity['x'])
                    
                    # Create line chart with highlights
                    fig = px.line(df_copy, x=opportunity['x'], y=opportunity['y'],
                                title=opportunity['title'],
                                labels={opportunity['y']: opportunity['y'].replace('_', ' ').title()})
                    
                    # Add high/low point highlights with friendly annotations
                    max_idx = df_copy[opportunity['y']].idxmax()
                    min_idx = df_copy[opportunity['y']].idxmin()
                    max_val = df_copy[opportunity['y']].max()
                    min_val = df_copy[opportunity['y']].min()
                    max_date = df_copy.loc[max_idx, opportunity['x']]
                    min_date = df_copy.loc[min_idx, opportunity['x']]
                    
                    # Add friendly annotations for high/low points
                    fig.add_annotation(
                        x=max_date, y=max_val,
                        text=f"{random.choice(annotations['high_point'])} {max_val:,.0f}",
                        showarrow=True, arrowcolor="green", bgcolor="lightgreen",
                        font=dict(color="darkgreen")
                    )
                    
                    fig.add_annotation(
                        x=min_date, y=min_val,
                        text=f"{random.choice(annotations['low_point'])} {min_val:,.0f}",
                        showarrow=True, arrowcolor="red", bgcolor="lightcoral",
                        font=dict(color="darkred")
                    )
                    
                    # Simple trend line if scipy available
                    if HAS_SCIPY and len(df_copy) <= 200:  # Even more restrictive for speed
                        try:
                            x_numeric = range(len(df_copy))
                            slope, intercept, r_value, _, _ = scipy.stats.linregress(x_numeric, df_copy[opportunity['y']])
                            trend_line = [slope * x + intercept for x in x_numeric]
                            
                            trend_text = "📈 Trending up" if slope > 0 else "📉 Trending down"
                            fig.add_trace(go.Scatter(
                                x=df_copy[opportunity['x']], y=trend_line,
                                mode='lines', name=f'{trend_text} (R²={r_value**2:.2f})',
                                line=dict(color='orange', width=2, dash='dash')
                            ))
                        except:
                            pass  # Skip trend if error
                    
                    fig.update_layout(template='plotly_white', showlegend=True)
                    return fig
                    
                except Exception:
                    return SmartChartGenerator._create_fallback_chart(data_to_use, opportunity)
                
            elif opportunity['type'] == 'kpi_metric':
                # Fast histogram with mean highlight
                fig = px.histogram(data_to_use, x=opportunity['column'], 
                                 title=opportunity['title'], nbins=15,
                                 labels={opportunity['column']: opportunity['column'].replace('_', ' ').title()})
                
                # Add mean line with friendly text
                mean_val = data_to_use[opportunity['column']].mean()
                median_val = data_to_use[opportunity['column']].median()
                
                fig.add_vline(x=mean_val, line_dash="dash", line_color="red", 
                             annotation_text=f"📊 Your average: {mean_val:.1f}")
                
                fig.update_layout(template='plotly_white', showlegend=False)
                return fig
                
            elif opportunity['type'] == 'category_performance':
                # Fast bar chart with top/bottom highlights
                top_categories = data_to_use[opportunity['x']].value_counts().head(8).index  # Reduced to 8 for speed
                df_filtered = data_to_use[data_to_use[opportunity['x']].isin(top_categories)]
                df_grouped = df_filtered.groupby(opportunity['x'])[opportunity['y']].mean().reset_index()
                df_grouped = df_grouped.sort_values(opportunity['y'], ascending=False)
                
                # Color-code bars (green for best, red for worst)
                colors = []
                for i, _ in enumerate(df_grouped.index):
                    if i == 0:  # Best performer
                        colors.append('lightgreen')
                    elif i == len(df_grouped) - 1:  # Worst performer  
                        colors.append('lightcoral')
                    else:
                        colors.append('lightblue')
                
                fig = px.bar(df_grouped, x=opportunity['x'], y=opportunity['y'],
                           title=opportunity['title'],
                           labels={opportunity['y']: f"Avg {opportunity['y'].replace('_', ' ').title()}"})
                
                # Update colors
                fig.update_traces(marker_color=colors)
                
                # Add value labels on bars
                fig.update_traces(texttemplate='%{y:.1f}', textposition='outside')
                
                # Add friendly annotations for best/worst
                best_val = df_grouped.iloc[0]
                worst_val = df_grouped.iloc[-1]
                
                fig.add_annotation(
                    x=best_val[opportunity['x']], y=best_val[opportunity['y']],
                    text=random.choice(annotations['high_point']), showarrow=False, 
                    font=dict(color="darkgreen", size=12)
                )
                
                if len(df_grouped) > 1:
                    fig.add_annotation(
                        x=worst_val[opportunity['x']], y=worst_val[opportunity['y']],
                        text=random.choice(annotations['low_point']), showarrow=False,
                        font=dict(color="darkred", size=12)
                    )
                
                fig.update_layout(template='plotly_white', xaxis_tickangle=-45, showlegend=False)
                return fig
                
            elif opportunity['type'] == 'correlation':
                # Super fast correlation heatmap
                corr_matrix = data_to_use[opportunity['columns'][:6]].corr()  # Limit to 6 columns max
                
                fig = px.imshow(corr_matrix, 
                              title=opportunity['title'],
                              color_continuous_scale='RdBu_r',
                              aspect="auto")
                
                # Add correlation values as text (simplified)
                fig.update_traces(text=corr_matrix.round(2).values, texttemplate="%{text}")
                fig.update_layout(template='plotly_white')
                return fig
                
            elif opportunity['type'] == 'geographic':
                # Geographic/regional analysis
                top_locations = data_to_use[opportunity['x']].value_counts().head(15).index
                df_filtered = data_to_use[data_to_use[opportunity['x']].isin(top_locations)]
                df_grouped = df_filtered.groupby(opportunity['x'])[opportunity['y']].mean().reset_index()
                df_grouped = df_grouped.sort_values(opportunity['y'], ascending=False)
                
                # Color-code by performance
                colors = []
                for i, _ in enumerate(df_grouped.index):
                    if i == 0:  # Best performer
                        colors.append('#2E8B57')  # Green
                    elif i < 3:  # Top 3
                        colors.append('#4682B4')  # Blue
                    elif i >= len(df_grouped) - 3:  # Bottom 3
                        colors.append('#DC143C')  # Red
                    else:
                        colors.append('#4169E1')  # Royal blue
                
                fig = px.bar(df_grouped, x=opportunity['x'], y=opportunity['y'],
                           title=opportunity['title'],
                           labels={opportunity['y']: f"Avg {opportunity['y'].replace('_', ' ').title()}"})
                
                # Update colors
                fig.update_traces(marker_color=colors)
                
                # Add value labels
                fig.update_traces(texttemplate='%{y:.1f}', textposition='outside')
                
                # Add friendly annotations for best/worst
                best_val = df_grouped.iloc[0]
                fig.add_annotation(
                    x=best_val[opportunity['x']], y=best_val[opportunity['y']],
                    text=random.choice(annotations['high_point']), showarrow=False, font=dict(size=16)
                )
                
                if len(df_grouped) > 1:
                    worst_val = df_grouped.iloc[-1]
                    fig.add_annotation(
                        x=worst_val[opportunity['x']], y=worst_val[opportunity['y']],
                        text=random.choice(annotations['low_point']), showarrow=False, font=dict(size=16)
                    )
                
                fig.update_layout(template='plotly_white', xaxis_tickangle=-45, showlegend=False)
                return fig
            
            elif opportunity['type'] == 'category_count':
                # Simple category count chart (fallback chart type)
                try:
                    value_counts = data_to_use[opportunity['column']].value_counts().head(10)
                    fig = px.bar(x=value_counts.index, y=value_counts.values,
                               title=opportunity['title'],
                               labels={'x': opportunity['column'].replace('_', ' ').title(), 'y': 'Count'})
                    
                    fig.update_layout(template='plotly_white', showlegend=False)
                    return fig
                except Exception:
                    return SmartChartGenerator._create_fallback_chart(data_to_use, opportunity)
                
        except Exception as e:
            st.warning(f"⚠️ Chart generation issue (using fallback): {str(e)}")
            return SmartChartGenerator._create_fallback_chart(data_to_use, opportunity)
        
        return None
    
    @staticmethod
    def _create_fallback_chart(df, opportunity):
        """Super simple fallback charts"""
        try:
            if opportunity['type'] == 'time_series':
                fig = px.scatter(df.sample(min(100, len(df))), 
                               x=opportunity['x'], y=opportunity['y'],
                               title=f"Sample: {opportunity['title']}")
                return fig
                
            elif opportunity['type'] == 'kpi_metric':
                fig = px.histogram(df.sample(min(100, len(df))), 
                                 x=opportunity['column'], title=f"Sample: {opportunity['title']}")
                return fig
                
            elif opportunity['type'] == 'category_performance':
                sample_df = df.sample(min(100, len(df)))
                fig = px.bar(sample_df.value_counts(opportunity['x']).head(5).reset_index(),
                           x='index', y=opportunity['x'], title=f"Top Categories: {opportunity['title']}")
                return fig
                
        except Exception:
            # Ultimate fallback - simple bar chart of value counts
            try:
                numeric_cols = df.select_dtypes(include=['number']).columns
                if len(numeric_cols) > 0:
                    col = numeric_cols[0]
                    fig = px.histogram(df.sample(min(50, len(df))), x=col, title="Data Sample")
                    return fig
            except:
                pass
        
        return None

class LLMInsights:
    """Generate AI-powered insights using Ollama (open source)"""
    
    @staticmethod
    def analyze_data_with_llm(df, chart_info, ollama_model='llama3.1', ollama_url='http://localhost:11434'):
        """Generate insights using Ollama with timeout"""
        
        # Create meaningful data context for LLM
        numeric_cols = df.select_dtypes(include=['number']).columns.tolist()
        categorical_cols = df.select_dtypes(include=['object']).columns.tolist()
        
        # Get actual data patterns for better insights
        insights_context = LLMInsights._build_data_context(df, chart_info)
        
        prompt = f"""You are a friendly data analyst helping someone understand their business data. Speak conversationally, as if explaining to a colleague.

{insights_context}

Chart: {chart_info.get('title', 'Data Analysis')}

Write 2-3 clear, conversational insights about what this data shows. Use natural language like "Here's what caught my eye" or "Something interesting I noticed". Avoid corporate jargon. Be specific with numbers when possible.

Focus on:
- What patterns you see that are worth mentioning
- Practical takeaways they can act on
- Interesting observations about their data

Write as if you're having a friendly conversation about their data.

Example tone:
"Here's what caught my eye: your sales really took off in March, jumping 23% compared to February. That's the kind of momentum that gets people excited! The data also shows your East region is consistently outperforming others."
"""
        
        try:
            # Try Ollama API with short timeout
            response = requests.post(
                f"{ollama_url}/api/generate",
                json={
                    'model': ollama_model,
                    'prompt': prompt,
                    'stream': False,
                    'options': {
                        'temperature': 0.7,
                        'top_p': 0.9,
                        'num_predict': 150  # Reduced for faster response
                    }
                },
                timeout=8  # Very short timeout for insights
            )
            
            if response.status_code == 200:
                result = response.json()
                return result.get('response', '').strip()
            
        except (requests.exceptions.RequestException, requests.exceptions.Timeout):
            # Quick fallback - don't wait for Ollama
            pass
        except Exception as e:
            # Any other error, use fallback
            pass
        
        # Enhanced fallback insights (always available)
        return LLMInsights.generate_business_insights(df, chart_info)
    
    @staticmethod
    def _build_data_context(df, chart_info):
        """Build meaningful context about the data for LLM"""
        context_parts = []
        
        # Dataset basics
        context_parts.append(f"Dataset: {df.shape[0]:,} records with {df.shape[1]} variables")
        
        # Key metrics summary
        numeric_cols = df.select_dtypes(include=['number']).columns
        if len(numeric_cols) > 0:
            key_stats = []
            for col in numeric_cols[:3]:  # Top 3 numeric columns
                mean_val = df[col].mean()
                median_val = df[col].median()
                if mean_val > 1000:
                    key_stats.append(f"{col}: avg ${mean_val:,.0f}")
                elif mean_val > 1:
                    key_stats.append(f"{col}: avg {mean_val:.1f}")
                else:
                    key_stats.append(f"{col}: avg {mean_val:.3f}")
            
            if key_stats:
                context_parts.append("Key metrics: " + ", ".join(key_stats))
        
        # Category insights
        categorical_cols = df.select_dtypes(include=['object']).columns
        if len(categorical_cols) > 0:
            cat_insights = []
            for col in categorical_cols[:2]:
                top_cat = df[col].value_counts().index[0]
                top_pct = (df[col].value_counts().iloc[0] / len(df)) * 100
                cat_insights.append(f"{col} dominated by {top_cat} ({top_pct:.0f}%)")
            
            if cat_insights:
                context_parts.append(" • ".join(cat_insights))
        
        return "\n".join(context_parts)
    
    @staticmethod
    def generate_business_insights(df, chart_info):
        """Generate human, conversational insights without LLM"""
        insights = []
        
        chart_title = chart_info.get('title', '').lower()
        numeric_cols = df.select_dtypes(include=['number']).columns
        categorical_cols = df.select_dtypes(include=['object']).columns
        
        # More conversational starters
        conversation_starters = [
            "Here's what caught my eye: ",
            "Something interesting I noticed: ",
            "The data tells an interesting story - ",
            "What stands out to me is ",
            "Looking at this data, ",
            "I found something worth mentioning: "
        ]
        
        # Time series insights - more human and specific
        if 'trend' in chart_title or 'time' in chart_title or 'journey' in chart_title or 'timeline' in chart_title:
            if len(numeric_cols) > 0:
                col = numeric_cols[0]
                try:
                    # Simple trend analysis
                    first_half = df[col][:len(df)//2].mean()
                    second_half = df[col][len(df)//2:].mean()
                    change_pct = ((second_half - first_half) / first_half) * 100
                    
                    starter = random.choice(conversation_starters)
                    col_name = col.replace('_', ' ').lower()
                    
                    if change_pct > 20:
                        insights.append(f"{starter}your {col_name} has really taken off, jumping {change_pct:.0f}% from the first half to the second half of your data. That's the kind of growth that gets people excited!")
                    elif change_pct > 10:
                        insights.append(f"{starter}there's solid upward momentum in your {col_name} - up {change_pct:.0f}% over time. Nothing too dramatic, but definitely heading in the right direction.")
                    elif change_pct < -20:
                        insights.append(f"{starter}your {col_name} has dropped quite a bit - down {abs(change_pct):.0f}%. That's definitely something worth digging into to understand what changed.")
                    elif change_pct < -10:
                        insights.append(f"{starter}there's been a {abs(change_pct):.0f}% dip in {col_name}. Not catastrophic, but worth keeping an eye on to see if it's just a temporary blip.")
                    else:
                        insights.append(f"{starter}your {col_name} has been pretty steady over time. Sometimes boring is good - it means you can count on consistent performance.")
                except:
                    insights.append(f"Your {chart_title} shows some interesting patterns over time that are worth exploring further.")
        
        # Performance by category insights - more relatable
        elif ('by' in chart_title or 'performs' in chart_title or 'breakdown' in chart_title) and len(categorical_cols) > 0 and len(numeric_cols) > 0:
            cat_col = categorical_cols[0]
            num_col = numeric_cols[0]
            
            try:
                performance = df.groupby(cat_col)[num_col].mean().sort_values(ascending=False)
                top_performer = performance.index[0]
                bottom_performer = performance.index[-1]
                performance_gap = ((performance.iloc[0] - performance.iloc[-1]) / performance.iloc[-1]) * 100
                
                col_name = num_col.replace('_', ' ').lower()
                cat_name = cat_col.replace('_', ' ').lower()
                
                starter = random.choice(conversation_starters)
                
                if performance_gap > 100:
                    insights.append(f"{starter}{top_performer} is absolutely crushing it compared to {bottom_performer} - we're talking {performance_gap:.0f}% better performance in {col_name}. Whatever they're doing, you might want to bottle it!")
                elif performance_gap > 50:
                    insights.append(f"{starter}{top_performer} is your clear winner, outperforming {bottom_performer} by {performance_gap:.0f}% in {col_name}. There's probably some good lessons to learn from their approach.")
                elif performance_gap > 20:
                    insights.append(f"{starter}while {top_performer} leads in {col_name}, the gap with {bottom_performer} isn't huge ({performance_gap:.0f}%). Most of your {cat_name} categories are performing reasonably well.")
                else:
                    insights.append(f"{starter}your {cat_name} categories show pretty similar performance in {col_name}. That consistency can be a good thing - no major weak spots to worry about.")
                
                # Check for consistency with more human language
                cv = performance.std() / performance.mean()
                if cv > 0.5:
                    insights.append(f"I noticed there's quite a bit of variation between different {cat_name} groups. Some are doing much better than others, which might mean there's room to help the lower performers catch up.")
                elif cv > 0.3:
                    insights.append(f"Performance across {cat_name} categories is moderately varied - not wildly different, but enough that you might want to understand what's driving the differences.")
                
            except:
                insights.append(f"There are some clear performance differences across your {cat_col.replace('_', ' ').lower()} categories that could be interesting to explore.")
        
        # Distribution insights - more accessible language
        elif 'distribution' in chart_title or 'spread' in chart_title or 'story behind' in chart_title:
            if len(numeric_cols) > 0:
                col = numeric_cols[0]
                try:
                    median = df[col].median()
                    mean = df[col].mean()
                    col_name = col.replace('_', ' ').lower()
                    
                    starter = random.choice(conversation_starters)
                    
                    # Skewness check with human language
                    if mean > median * 1.3:
                        high_val = df[col].max()
                        insights.append(f"{starter}most of your {col_name} values are pretty typical, but you've got some real standouts pulling the average up. Your highest value is {high_val:,.0f}, which is way above what most entries show.")
                    elif mean < median * 0.7:
                        insights.append(f"{starter}you have a few lower values in {col_name} that are dragging down your overall average. The middle value is actually higher than the average, which usually means most of your data is performing better than the average suggests.")
                    else:
                        insights.append(f"{starter}your {col_name} values are spread out pretty evenly - no major outliers throwing things off. That's often a sign of consistent, predictable performance.")
                        
                except:
                    col_name = col.replace('_', ' ').lower()
                    insights.append(f"The spread of your {col_name} values shows some interesting patterns that might help you understand what drives different outcomes.")
        
        # Correlation insights - simplified
        elif 'correlation' in chart_title or 'connect' in chart_title:
            if len(numeric_cols) >= 2:
                try:
                    corr_matrix = df[numeric_cols].corr()
                    strong_corrs = []
                    
                    for i in range(len(corr_matrix.columns)):
                        for j in range(i+1, len(corr_matrix.columns)):
                            corr_val = corr_matrix.iloc[i, j]
                            if abs(corr_val) > 0.7:
                                col1 = corr_matrix.columns[i].replace('_', ' ').lower()
                                col2 = corr_matrix.columns[j].replace('_', ' ').lower()
                                
                                if corr_val > 0:
                                    strong_corrs.append(f"when {col1} goes up, {col2} tends to go up too")
                                else:
                                    strong_corrs.append(f"when {col1} increases, {col2} usually decreases")
                    
                    starter = random.choice(conversation_starters)
                    
                    if strong_corrs:
                        insights.append(f"{starter}{strong_corrs[0]}. That kind of relationship can be really useful for predicting what might happen next.")
                    else:
                        insights.append(f"{starter}your different metrics seem to dance to their own beat - they're not strongly connected to each other. That means each one is probably measuring something unique and valuable.")
                        
                except:
                    insights.append("The relationships between your different metrics reveal some interesting patterns about how your data behaves.")
        
        # Geographic insights - more engaging
        elif 'geographic' in chart_title.lower() or any(geo_word in chart_title.lower() for geo_word in ['country', 'state', 'city', 'region']):
            if len(categorical_cols) > 0 and len(numeric_cols) > 0:
                geo_col = categorical_cols[0]
                num_col = numeric_cols[0]
                
                try:
                    performance = df.groupby(geo_col)[num_col].mean().sort_values(ascending=False)
                    top_location = performance.index[0]
                    bottom_location = performance.index[-1]
                    
                    geo_name = geo_col.replace('_', ' ').lower()
                    metric_name = num_col.replace('_', ' ').lower()
                    
                    starter = random.choice(conversation_starters)
                    insights.append(f"{starter}{top_location} is your star performer when it comes to {metric_name}, while {bottom_location} might need some extra attention. Geographic differences like this often tell a story about local market conditions or operational differences.")
                    
                except:
                    insights.append(f"There are some interesting geographic patterns in your data that might reflect local market conditions or regional differences.")
        
        # Default insights - more personal and specific
        if len(insights) == 0:
            data_size_insights = [
                f"You're working with {len(df):,} records here - that's a solid amount of data to draw meaningful conclusions from.",
                f"With {len(df):,} data points, you've got enough information to spot real patterns and trends.",
                f"This dataset of {len(df):,} records gives you a nice foundation for making data-driven decisions."
            ]
            insights.append(random.choice(data_size_insights))
            
            # Data quality insight - more conversational
            missing_pct = (df.isnull().sum().sum() / (df.shape[0] * df.shape[1])) * 100
            if missing_pct < 2:
                quality_insights = [
                    "Your data quality looks excellent - hardly any missing values to worry about.",
                    "I'm impressed by how complete your dataset is. Clean data like this makes analysis so much more reliable.",
                    "Data quality wise, you're in great shape. Almost no gaps or missing information."
                ]
            elif missing_pct < 10:
                quality_insights = [
                    f"Your data is mostly complete with only about {missing_pct:.1f}% missing values - that's pretty typical and nothing to worry about.",
                    f"Data quality looks good overall. You have some missing values ({missing_pct:.1f}%) but nothing that would throw off your analysis.",
                    f"With {missing_pct:.1f}% missing data, you're well within the normal range. Most datasets have some gaps."
                ]
            else:
                quality_insights = [
                    f"You have about {missing_pct:.1f}% missing values, which is noticeable but still workable for most analyses.",
                    f"There are some gaps in the data ({missing_pct:.1f}% missing), but that's not uncommon with real-world datasets."
                ]
            
            insights.append(random.choice(quality_insights))
        
        # Add a forward-looking insight
        if len(insights) < 3:
            action_insights = [
                "The patterns here definitely give you some directions to explore further.",
                "These insights are just the starting point - there's probably more interesting stuff hiding in the details.",
                "This gives you a good foundation, but I'd be curious to see what you discover when you dig deeper into specific areas.",
                "Based on what I'm seeing, you've got some clear next steps for investigation."
            ]
            insights.append(random.choice(action_insights))
        
        # Return up to 2 insights for readability, joined naturally
        selected_insights = insights[:2]
        return " ".join(selected_insights)

def get_basic_data_info(data, question):
    """Provide basic data information without AI - enhanced with personality"""
    question_lower = question.lower()
    
    # More conversational responses
    friendly_starters = [
        "Here's what I can tell you: ",
        "Let me break this down for you: ",
        "Good question! ",
        "I took a look and here's what I found: "
    ]
    
    # Basic data info
    numeric_cols = data.select_dtypes(include=['number']).columns.tolist()
    categorical_cols = data.select_dtypes(include=['object']).columns.tolist()
    
    # Marketing-specific analysis - more engaging
    if any(word in question_lower for word in ['marketing', 'campaign', 'ctr', 'conversion', 'performance']):
        starter = random.choice(friendly_starters)
        marketing_insights = []
        
        # CTR analysis with personality
        ctr_cols = [col for col in data.columns if 'ctr' in col.lower()]
        if ctr_cols:
            ctr_col = ctr_cols[0]
            avg_ctr = data[ctr_col].mean()
            max_ctr = data[ctr_col].max()
            min_ctr = data[ctr_col].min()
            
            if avg_ctr > 3:
                marketing_insights.append(f"Your average CTR of {avg_ctr:.2f}% is pretty solid! Your best performing piece hit {max_ctr:.2f}%.")
            elif avg_ctr > 1:
                marketing_insights.append(f"Your CTR averages {avg_ctr:.2f}% - not bad, though there's room to improve. Your top performer reached {max_ctr:.2f}%.")
            else:
                marketing_insights.append(f"Your CTR is averaging {avg_ctr:.2f}%, which suggests there might be room for optimization. Your best piece did {max_ctr:.2f}% though!")
        
        # Impressions analysis
        impression_cols = [col for col in data.columns if 'impression' in col.lower()]
        if impression_cols:
            imp_col = impression_cols[0]
            avg_imp = data[imp_col].mean()
            max_imp = data[imp_col].max()
            marketing_insights.append(f"You're averaging {avg_imp:,.0f} impressions, with your best piece getting {max_imp:,.0f}.")
        
        # Click analysis
        click_cols = [col for col in data.columns if 'click' in col.lower()]
        if click_cols:
            click_col = click_cols[0]
            avg_clicks = data[click_col].mean()
            total_clicks = data[click_col].sum()
            marketing_insights.append(f"You've gotten {total_clicks:,.0f} total clicks, averaging {avg_clicks:,.0f} per campaign.")
        
        # Campaign type analysis
        campaign_cols = [col for col in categorical_cols if any(word in col.lower() for word in ['campaign', 'type', 'channel', 'source'])]
        if campaign_cols:
            camp_col = campaign_cols[0]
            top_campaigns = data[camp_col].value_counts().head(3)
            campaign_list = ', '.join([f"{idx} ({count})" for idx, count in top_campaigns.items()])
            marketing_insights.append(f"Your top performing {camp_col.replace('_', ' ').lower()}s are: {campaign_list}.")
        
        if marketing_insights:
            return f"{starter}{chr(10).join(marketing_insights)}\n\nWant to see the visual breakdown? Try generating some charts!"
    
    if any(word in question_lower for word in ['columns', 'column', 'fields']):
        starter = random.choice(friendly_starters)
        return f"""{starter}You've got {len(data.columns)} columns to work with.

**Numbers to crunch ({len(numeric_cols)}):** {', '.join(numeric_cols[:8])}{'...' if len(numeric_cols) > 8 else ''}

**Categories to explore ({len(categorical_cols)}):** {', '.join(categorical_cols[:8])}{'...' if len(categorical_cols) > 8 else ''}

That's a nice mix to work with!"""

    elif any(word in question_lower for word in ['rows', 'size', 'shape', 'big', 'large']):
        missing = data.isnull().sum().sum()
        complete = ((data.shape[0] * data.shape[1]) - missing) / (data.shape[0] * data.shape[1]) * 100
        starter = random.choice(friendly_starters)
        return f"""{starter}Your dataset is a nice size to work with.

• **Rows:** {data.shape[0]:,} records
• **Columns:** {data.shape[1]} different variables
• **Total data points:** {data.shape[0] * data.shape[1]:,}
• **Data completeness:** {complete:.1f}% complete

{f"You have {missing:,} missing values scattered throughout, which is pretty normal." if missing > 0 else "Everything's filled in - no missing values at all!"}"""

    elif any(word in question_lower for word in ['stats', 'statistics', 'summary', 'describe']):
        if numeric_cols:
            first_col = numeric_cols[0]
            stats = data[first_col].describe()
            starter = random.choice(friendly_starters)
            return f"""{starter}Let me give you a quick rundown of your {first_col.replace('_', ' ').lower()} numbers:

• **Average:** {stats['mean']:.2f}
• **Typical value (median):** {stats['50%']:.2f}
• **Range:** {stats['min']:.2f} to {stats['max']:.2f}
• **How spread out:** {stats['std']:.2f} standard deviation

Want to see this visualized? Hit the chart generation button!"""

    elif any(word in question_lower for word in ['missing', 'null', 'empty', 'quality']):
        missing_by_col = data.isnull().sum()
        problematic = missing_by_col[missing_by_col > 0].head(5)
        starter = random.choice(friendly_starters)
        if len(problematic) > 0:
            missing_info = '\n'.join([f"• {col}: {count:,} missing ({count/len(data)*100:.1f}%)" for col, count in problematic.items()])
            return f"""{starter}Your data quality is pretty typical - here's what's missing:

{missing_info}

**Overall:** {((data.shape[0] * data.shape[1]) - data.isnull().sum().sum()) / (data.shape[0] * data.shape[1]) * 100:.1f}% of your data is complete, which is totally workable!"""
        else:
            return f"{starter}Your data quality is excellent! No missing values anywhere - that makes analysis much cleaner."

    elif any(word in question_lower for word in ['trend', 'pattern', 'insight', 'analysis']):
        # Enhanced trend analysis
        trend_insights = []
        starter = random.choice(friendly_starters)
        
        # Analyze numeric trends
        if numeric_cols:
            for col in numeric_cols[:3]:
                series = data[col].dropna()
                if len(series) > 10:
                    # Simple trend analysis
                    first_quarter = series.iloc[:len(series)//4].mean()
                    last_quarter = series.iloc[-len(series)//4:].mean()
                    change = ((last_quarter - first_quarter) / first_quarter) * 100 if first_quarter > 0 else 0
                    
                    col_name = col.replace('_', ' ').lower()
                    
                    if abs(change) > 20:
                        direction = "climbing up" if change > 0 else "dropping down"
                        trend_insights.append(f"Your {col_name} is {direction} quite a bit - about {abs(change):.1f}% change from beginning to end.")
                    elif abs(change) > 10:
                        direction = "trending upward" if change > 0 else "trending downward"
                        trend_insights.append(f"Your {col_name} is {direction} with a {abs(change):.1f}% shift overall.")
                    else:
                        trend_insights.append(f"Your {col_name} stays pretty consistent (just ±{abs(change):.1f}% variation).")
        
        # Analyze categorical patterns
        if categorical_cols:
            for col in categorical_cols[:2]:
                if data[col].nunique() <= 20:
                    top_cat = data[col].value_counts()
                    dominant = top_cat.iloc[0]
                    total = len(data)
                    dominance_pct = (dominant / total) * 100
                    col_name = col.replace('_', ' ').lower()
                    if dominance_pct > 50:
                        trend_insights.append(f"Your {col_name} data is dominated by {top_cat.index[0]} - it makes up {dominance_pct:.1f}% of everything.")
                    else:
                        trend_insights.append(f"Your {col_name} categories are pretty well distributed, with {top_cat.index[0]} leading at {dominance_pct:.1f}%.")
        
        # Marketing-specific insights
        marketing_keywords = ['ctr', 'click', 'impression', 'conversion', 'spend', 'cost']
        marketing_cols = [col for col in numeric_cols if any(kw in col.lower() for kw in marketing_keywords)]
        
        if marketing_cols:
            trend_insights.append(f"I can see you've got marketing metrics here: {', '.join([col.replace('_', ' ').lower() for col in marketing_cols[:3]])}.")
            
            # CTR analysis if available
            ctr_cols = [col for col in data.columns if 'ctr' in col.lower()]
            if ctr_cols:
                ctr_col = ctr_cols[0]
                avg_ctr = data[ctr_col].mean()
                high_ctr = data[ctr_col].quantile(0.9)
                trend_insights.append(f"Your CTR averages {avg_ctr:.2f}%, with the top 10% of campaigns hitting {high_ctr:.2f}%.")
        
        if trend_insights:
            insights_text = '\n'.join([f"• {insight}" for insight in trend_insights[:4]])  # Limit to top 4 insights
            return f"""{starter}I found some interesting patterns:

{insights_text}

Want to see these trends visualized? The chart generator above will create some really nice visuals for you!"""
        else:
            return f"""{starter}Your dataset has {data.shape[0]:,} records with {len(numeric_cols)} numeric and {len(categorical_cols)} text-based variables. 

To spot the interesting trends and patterns, try:
1. Hitting that "Generate Charts" button above
2. Looking for time-based patterns
3. Comparing different categories

Want to ask about something more specific? Try "What's my average CTR?" or "Which columns have the highest values?"."""

    else:
        starters = [
            "I'm here to help you explore your data! ",
            "Let's dig into your dataset! ",
            "Happy to help you understand your data better! "
        ]
        
        return f"""{random.choice(starters)}

**Your dataset:** {data.shape[0]:,} rows × {data.shape[1]} columns

**Try asking me:**
• "What's the story behind my marketing data?"
• "Which columns have the most interesting patterns?"
• "Help me understand my data quality"
• "What should I focus on first?"

**For visual insights:** Hit that chart generation button above - it'll create up to 8 different views of your data!"""

def chat_with_data(data, question, ollama_model='llama3.1', ollama_url='http://localhost:11434'):
    """Chat interface for data exploration using Ollama with fallback"""
    if data is None:
        return "Please load a dataset first to start chatting with your data."
    
    # Create comprehensive data context
    numeric_cols = data.select_dtypes(include=['number']).columns.tolist()
    categorical_cols = data.select_dtypes(include=['object']).columns.tolist()
    
    # Sample of actual data for context (smaller sample for speed)
    data_sample = data.head(2).to_string() if len(data) > 0 else "No data available"
    
    # Statistical summaries for numeric columns (limit for speed)
    numeric_summary = ""
    if numeric_cols:
        for col in numeric_cols[:3]:
            mean_val = data[col].mean()
            max_val = data[col].max()
            min_val = data[col].min()
            numeric_summary += f"• {col}: range {min_val:.1f} to {max_val:.1f}, average {mean_val:.1f}\n"
    
    # Category summaries (limit for speed)
    category_summary = ""
    if categorical_cols:
        for col in categorical_cols[:2]:
            top_categories = data[col].value_counts().head(2)
            category_summary += f"• {col}: top values are {', '.join([f'{idx} ({count})' for idx, count in top_categories.items()])}\n"
    
    data_context = f"""You are a friendly data analyst helping someone understand their business data. Speak conversationally and naturally, like you're talking to a colleague over coffee.

Dataset: {data.shape[0]} rows and {data.shape[1]} columns

Key Numbers:
{numeric_summary[:200]}

Main Categories:
{category_summary[:200]}

User Question: {question}

Respond in a conversational, helpful way. Use phrases like "Here's what I found" or "Something interesting I noticed" rather than formal business language. Be specific with numbers when you can. Keep it friendly and practical."""
    
    try:
        # Try Ollama API with longer timeout for analytical questions
        timeout_duration = 15 if any(word in question.lower() for word in ['trend', 'pattern', 'insight', 'analysis', 'correlation']) else 8
        
        response = requests.post(
            f"{ollama_url}/api/generate",
            json={
                'model': ollama_model,
                'prompt': data_context,
                'stream': False,
                'options': {
                    'temperature': 0.3,
                    'top_p': 0.9,
                    'num_predict': 150
                }
            },
            timeout=timeout_duration  # Longer timeout for analytical questions
        )
        
        if response.status_code == 200:
            result = response.json()
            ai_response = result.get('response', '').strip()
            if ai_response and len(ai_response) > 10:  # Valid response
                return ai_response
            else:
                # Empty response, use fallback
                return get_basic_data_info(data, question)
        else:
            # Bad status code, use fallback
            return get_basic_data_info(data, question)
            
    except requests.exceptions.ConnectionError:
        # Ollama not running - provide helpful info
        basic_info = get_basic_data_info(data, question)
        return f"""{basic_info}

🤖 **Want AI-powered insights?**
1. Install Ollama: https://ollama.ai
2. Run: `ollama serve`
3. Download model: `ollama pull llama3.1`"""
        
    except requests.exceptions.Timeout:
        # Timeout - provide enhanced basic response for analytical questions
        if any(word in question.lower() for word in ['trend', 'pattern', 'insight', 'analysis']):
            return get_basic_data_info(data, question)  # This now has enhanced trend analysis
        else:
            return f"""⏰ AI response timed out. Here's what I can tell you:

{get_basic_data_info(data, question)}

**Tip:** For instant responses, try questions like "What columns are available?" or "Show me key statistics"."""
        
    except Exception as e:
        # Any other error - provide basic analysis
        return get_basic_data_info(data, question)

class KPIDashboard:
    def __init__(self):
        self.data = None
        self.charts = {}
        self.data_source = None
        
    def reset_data(self):
        """Reset all data and charts"""
        self.data = None
        self.charts = {}
        self.data_source = None
        
    def load_csv_upload(self, uploaded_file):
        """Load CSV from uploaded file"""
        try:
            self.data = pd.read_csv(uploaded_file)
            self.data_source = uploaded_file.name
            return True, f"Successfully loaded {uploaded_file.name}"
        except Exception as e:
            return False, f"Error loading CSV: {str(e)}"
        
    def load_kaggle_dataset(self, dataset_name):
        """Load dataset from Kaggle"""
        try:
            # Reset previous data
            self.reset_data()
            
            # Download dataset
            kaggle.api.dataset_download_files(dataset_name, path='./data', unzip=True)
            
            # Find CSV files in the data directory
            data_files = [f for f in os.listdir('./data') if f.endswith('.csv')]
            
            if data_files:
                # Load the first CSV file found
                file_path = f"./data/{data_files[0]}"
                self.data = pd.read_csv(file_path)
                self.data_source = f"Kaggle: {dataset_name}"
                return True, f"Successfully loaded {data_files[0]} from Kaggle"
            else:
                return False, "No CSV files found in the dataset"
                
        except Exception as e:
            return False, f"Error loading dataset: {str(e)}"
    
    def generate_smart_charts(self, use_full_dataset=None, sample_size=None):
        """Generate intelligent, business-relevant charts with user preferences"""
        if self.data is None:
            return {}
        
        # Use defaults if not provided
        if use_full_dataset is None:
            use_full_dataset = len(self.data) <= 5000
        if sample_size is None:
            sample_size = 2000
        
        # Reset charts
        self.charts = {}
        
        # Get chart opportunities with user preferences using the global function
        opportunities = identify_chart_opportunities(
            self.data, use_full_dataset, sample_size
        )
        
        if not opportunities:
            st.warning("No suitable chart opportunities found in this dataset.")
            return {}
        
        # Create charts without progress bars for maximum speed
        for i, opportunity in enumerate(opportunities):
            try:
                chart = SmartChartGenerator.create_chart(self.data, opportunity)
                if chart:
                    self.charts[f'chart_{i+1}'] = {
                        'title': opportunity['title'],
                        'figure': chart,
                        'opportunity': opportunity
                    }
            except Exception as e:
                st.error(f"❌ Error creating {opportunity['title']}: {str(e)}")
        
        return self.charts
    
    def create_powerpoint_with_charts(self, charts, insights):
        """Create PowerPoint with actual chart images"""
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "KPI Analysis Report"
        subtitle.text = f"Data Source: {self.data_source}\nGenerated: {datetime.now().strftime('%B %d, %Y')}"
        
        # Executive Summary slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Executive Summary"
        tf = content.text_frame
        tf.clear()
        
        # Add data overview
        p = tf.paragraphs[0]
        p.text = f"Dataset Overview: {self.data.shape[0]:,} rows, {self.data.shape[1]} columns"
        
        # Add key insights
        for insight in insights[:4]:
            p = tf.add_paragraph()
            p.text = f"• {insight['insight'][:100]}..."
            p.level = 0
        
        # Individual chart slides
        for chart_name, chart_info in charts.items():
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
            
            # Add title
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_shape.text_frame
            title_frame.text = chart_info['title']
            title_frame.paragraphs[0].font.size = Inches(0.3)
            title_frame.paragraphs[0].font.bold = True
            
            # Save chart as image and add to slide
            try:
                # Save plotly chart as image
                img_bytes = chart_info['figure'].to_image(format="png", width=800, height=500)
                
                # Create temporary file
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                    tmp_file.write(img_bytes)
                    tmp_path = tmp_file.name
                
                # Add image to slide
                slide.shapes.add_picture(tmp_path, Inches(1), Inches(1.5), Inches(8), Inches(5))
                
                # Clean up temp file
                os.unlink(tmp_path)
                
                # Add insight text
                insight_text = next((insight['insight'] for insight in insights if insight['chart'] == chart_name), 
                                  "Key business insights can be derived from this visualization.")
                
                text_box = slide.shapes.add_textbox(Inches(1), Inches(7), Inches(8), Inches(1))
                text_frame = text_box.text_frame
                text_frame.text = f"Insight: {insight_text}"
                text_frame.word_wrap = True
                
            except Exception as e:
                # Fallback: add text-only slide
                content_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
                content_frame = content_shape.text_frame
                content_frame.text = f"Chart: {chart_info['title']}\n\nNote: Chart image could not be embedded.\nError: {str(e)}"
        
        # Save presentation
        os.makedirs('generated_ppts', exist_ok=True)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"generated_ppts/KPI_Report_{timestamp}.pptx"
        prs.save(filename)
        
        return filename

def main():
    st.title("📊 KPI Insights Dashboard v2.5")
    st.markdown("**Human-Like Insights Version** - Now with conversational, friendly analysis that actually makes sense!")
    
    # Initialize dashboard
    if 'dashboard' not in st.session_state:
        st.session_state.dashboard = KPIDashboard()
    
    # Sidebar
    st.sidebar.title("🚀 Dashboard Controls")
    
    # Step 1: Data Loading
    st.sidebar.header("1. Load Data")
    
    # Data source selection
    data_source = st.sidebar.radio(
        "Choose data source:",
        ["Upload CSV", "Kaggle Dataset"]
    )
    
    if data_source == "Upload CSV":
        uploaded_file = st.sidebar.file_uploader("Choose a CSV file", type="csv")
        if uploaded_file is not None:
            if st.sidebar.button("Load CSV"):
                with st.spinner("Loading CSV file..."):
                    success, message = st.session_state.dashboard.load_csv_upload(uploaded_file)
                    if success:
                        st.sidebar.success(message)
                        st.rerun()
                    else:
                        st.sidebar.error(message)
    
    else:  # Kaggle Dataset
        dataset_name = st.sidebar.text_input(
            "Enter Kaggle Dataset Name",
            placeholder="e.g., username/dataset-name",
            help="Format: username/dataset-name (as shown in Kaggle URL)"
        )
        
        if st.sidebar.button("Load Kaggle Dataset"):
            if dataset_name:
                with st.spinner("Loading dataset from Kaggle..."):
                    success, message = st.session_state.dashboard.load_kaggle_dataset(dataset_name)
                    if success:
                        st.sidebar.success(message)
                        st.rerun()
                    else:
                        st.sidebar.error(message)
            else:
                st.sidebar.error("Please enter a dataset name")
        
        # Sample datasets for testing
        st.sidebar.subheader("Sample Datasets:")
        sample_datasets = [
            "russellyates88/stock-market-data",
            "prasadperera/the-boston-housing-dataset",
            "vikramtiwari/pima-indians-diabetes-database"
        ]
        
        for dataset in sample_datasets:
            if st.sidebar.button(f"📊 {dataset.split('/')[1][:20]}...", key=dataset):
                with st.spinner(f"Loading {dataset}..."):
                    success, message = st.session_state.dashboard.load_kaggle_dataset(dataset)
                    if success:
                        st.sidebar.success(message)
                        st.rerun()
                    else:
                        st.sidebar.error(message)
    
    # Ollama Configuration
    st.sidebar.header("2. AI Configuration")
    
    with st.sidebar.expander("🤖 Ollama Setup", expanded=False):
        st.markdown("""
        **For Enhanced AI Insights:**
        1. Install Ollama: https://ollama.ai
        2. Run: `ollama serve`
        3. Download model: `ollama pull llama3.1`
        """)
    
    ollama_url = st.sidebar.text_input("Ollama URL", value="http://localhost:11434", 
                                      help="Default: http://localhost:11434")
    st.session_state.ollama_url = ollama_url
    
    ollama_model = st.sidebar.selectbox("Ollama Model", 
                                       ["llama3.1", "llama3", "mistral", "codellama"],
                                       help="Make sure the model is downloaded: ollama pull <model>")
    st.session_state.ollama_model = ollama_model
    
    # Analysis Options
    st.sidebar.header("3. Analysis Options")
    
    # Set initial values
    if st.session_state.dashboard.data is not None:
        default_full_dataset = len(st.session_state.dashboard.data) <= 5000
        max_sample = min(5000, len(st.session_state.dashboard.data))
    else:
        default_full_dataset = True
        max_sample = 5000
    
    use_full_dataset = st.sidebar.checkbox(
        "Use Full Dataset", 
        value=default_full_dataset,
        help="Uncheck for faster analysis with sampling"
    )
    st.session_state.use_full_dataset = use_full_dataset
    
    if not use_full_dataset and st.session_state.dashboard.data is not None:
        sample_size = st.sidebar.slider(
            "Sample Size", 
            min_value=500, 
            max_value=max_sample, 
            value=min(2000, max_sample),
            help="Smaller = Faster, Larger = More accurate"
        )
        st.session_state.sample_size = sample_size
        st.sidebar.info(f"Will analyze {sample_size:,} out of {len(st.session_state.dashboard.data):,} rows")
    else:
        st.session_state.sample_size = 2000  # Default when using full dataset
    
    # Test Ollama connection
    if st.sidebar.button("🔌 Test Ollama Connection"):
        try:
            response = requests.get(f"{ollama_url}/api/tags", timeout=5)
            if response.status_code == 200:
                models = response.json().get('models', [])
                if models:
                    available_models = [m['name'] for m in models]
                    st.sidebar.success(f"✅ Connected! Available models: {', '.join(available_models[:3])}")
                else:
                    st.sidebar.warning("Ollama connected but no models found. Run: ollama pull llama3.1")
            else:
                st.sidebar.error("❌ Ollama not responding")
        except:
            st.sidebar.error("❌ Cannot connect to Ollama. Make sure it's running!")
            st.sidebar.info("Start with: `ollama serve`")
    
    # Main content area
    if st.session_state.dashboard.data is not None:
        data = st.session_state.dashboard.data
        
        # Data preview
        st.header("📋 Data Overview")
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            st.metric("📊 Total Rows", f"{len(data):,}")
        with col2:
            st.metric("📈 Total Columns", len(data.columns))
        with col3:
            st.metric("❌ Missing Values", f"{data.isnull().sum().sum():,}")
        with col4:
            st.metric("💾 Memory Usage", f"{data.memory_usage(deep=True).sum() / 1024**2:.1f} MB")
        
        # Data source info
        st.info(f"**Data Source:** {st.session_state.dashboard.data_source}")
        
        # Show data sample
        with st.expander("📊 View Data Sample", expanded=False):
            st.dataframe(data.head(100), use_container_width=True)
        
        # Generate charts with better feedback
        if st.button("🚀 Generate Charts & Human-Like Insights", type="primary"):
            # Quick data validation
            if len(data) == 0:
                st.error("❌ Dataset is empty! Please load a different dataset.")
                return
            
            if len(data.columns) < 2:
                st.error("❌ Dataset needs at least 2 columns for meaningful analysis.")
                return
            
            # Show what we're working with
            numeric_cols = data.select_dtypes(include=['number']).columns.tolist()
            categorical_cols = data.select_dtypes(include=['object']).columns.tolist()
            
            with st.spinner(random.choice(friendly_progress_messages)):
                # Get user preferences from sidebar (with fallbacks)
                try:
                    use_full = st.session_state.get('use_full_dataset', len(data) <= 5000)
                    sample_sz = st.session_state.get('sample_size', 2000)
                except:
                    use_full = len(data) <= 5000
                    sample_sz = 2000
                
                # Show dataset info based on user preferences
                if use_full or len(data) <= 5000:
                    analysis_scope = "Full dataset"
                    analysis_size = len(data)
                else:
                    analysis_scope = f"Sample of {sample_sz:,} rows"
                    analysis_size = min(sample_sz, len(data))
                
                st.info(f"""
                ⚡ **What I'm Creating:**
                - Dataset size: {len(data):,} rows
                - Analysis scope: {analysis_scope}
                - Processing: {analysis_size:,} rows for visualization
                - Charts: Up to 8 engaging visuals with friendly insights
                - Style: Conversational analysis that actually makes sense!
                """)
                
                charts = st.session_state.dashboard.generate_smart_charts(
                    use_full_dataset=use_full,
                    sample_size=sample_sz
                )
            
            if charts:
                # Add option to skip AI insights for faster generation
                col1, col2 = st.columns([3, 1])
                with col1:
                    st.success(f"✅ Generated {len(charts)} charts with human-like insights!")
                with col2:
                    skip_insights = st.button("⚡ Skip AI Enhancement", help="Show charts with basic insights for speed")
                
                # Generate insights based on user choice
                if skip_insights:
                    # Create friendly insights without AI
                    insights = []
                    for chart_name, chart_info in charts.items():
                        basic_insight = LLMInsights.generate_business_insights(data, chart_info)
                        insights.append({
                            'chart': chart_name,
                            'title': chart_info['title'],
                            'insight': basic_insight
                        })
                    st.info("⚡ Charts ready with friendly insights! AI enhancement skipped for maximum speed.")
                else:
                    # Generate AI insights with timeout
                    with st.spinner("🤖 Creating human-like insights (max 15 seconds)..."):
                        insights = []
                        insight_timeout = 15  # 15 seconds total for all insights
                        start_time = datetime.now()
                        
                        for chart_name, chart_info in charts.items():
                            if (datetime.now() - start_time).seconds > insight_timeout:
                                # Timeout reached, use friendly fallback insights
                                friendly_insight = LLMInsights.generate_business_insights(data, chart_info)
                                insights.append({
                                    'chart': chart_name,
                                    'title': chart_info['title'],
                                    'insight': friendly_insight
                                })
                                continue
                                
                            try:
                                insight_text = LLMInsights.analyze_data_with_llm(
                                    data, chart_info, ollama_model, ollama_url
                                )
                                insights.append({
                                    'chart': chart_name,
                                    'title': chart_info['title'],
                                    'insight': insight_text
                                })
                            except Exception:
                                # Fallback to friendly human insights
                                friendly_insight = LLMInsights.generate_business_insights(data, chart_info)
                                insights.append({
                                    'chart': chart_name,
                                    'title': chart_info['title'],
                                    'insight': friendly_insight
                                })
                
                # Store in session state
                st.session_state.charts = charts
                st.session_state.insights = insights
                
                st.balloons()  # Celebrate success!
            else:
                st.error("❌ Could not generate charts. The dataset format may not be suitable for automatic analysis.")
        
        # Display charts
        if 'charts' in st.session_state and st.session_state.charts:
            st.header("📊 Your Data Stories")
            st.markdown("*Charts with conversational insights that actually make sense*")
            
            # Create tabs for better organization
            chart_names = list(st.session_state.charts.keys())
            if len(chart_names) > 1:
                tabs = st.tabs([st.session_state.charts[name]['title'] for name in chart_names])
                
                for i, (chart_name, tab) in enumerate(zip(chart_names, tabs)):
                    with tab:
                        chart_info = st.session_state.charts[chart_name]
                        st.plotly_chart(
                            chart_info['figure'], 
                            use_container_width=True
                        )
                        
                        # Show corresponding insight
                        if 'insights' in st.session_state:
                            insight = next((ins for ins in st.session_state.insights 
                                          if ins['chart'] == chart_name), None)
                            if insight:
                                st.info(f"💡 **What I Found:** {insight['insight']}")
            else:
                # Single chart
                for chart_name, chart_info in st.session_state.charts.items():
                    st.subheader(chart_info['title'])
                    st.plotly_chart(
                        chart_info['figure'], 
                        use_container_width=True
                    )
        
        # Generate PowerPoint
        if 'insights' in st.session_state and st.session_state.insights:
            st.header("💭 Friendly Data Insights")
            st.markdown("*No corporate jargon here - just clear, helpful observations about your data*")
            
            for insight in st.session_state.insights:
                with st.expander(f"💡 {insight['title']}", expanded=True):
                    st.write(insight['insight'])
            
            col1, col2 = st.columns([1, 1])
            with col1:
                if st.button("🎯 Generate PowerPoint Report", type="primary"):
                    with st.spinner("Creating PowerPoint with your insights..."):
                        try:
                            filename = st.session_state.dashboard.create_powerpoint_with_charts(
                                st.session_state.charts, 
                                st.session_state.insights
                            )
                            st.success(f"✅ PowerPoint created: {filename}")
                            
                            # Provide download link
                            with open(filename, "rb") as file:
                                st.download_button(
                                    label="📥 Download PowerPoint",
                                    data=file.read(),
                                    file_name=filename.split('/')[-1],
                                    mime="application/vnd.ms-powerpoint"
                                )
                        except Exception as e:
                            st.error(f"Error creating PowerPoint: {str(e)}")
        
        # Chat with Data Section  
        st.header("💬 Chat with Your Data")
        st.markdown("*Ask questions in plain English and get friendly, helpful answers*")
        
        if "chat_history" not in st.session_state:
            st.session_state.chat_history = []
        
        # Chat input
        user_question = st.chat_input("Ask me anything about your data...")
        
        if user_question:
            # Add user message to chat history
            st.session_state.chat_history.append({"role": "user", "content": user_question})
            
            # Show what's happening
            status_placeholder = st.empty()
            
            # Get AI response
            with st.spinner("🤖 Looking into that for you..."):
                try:
                    # Get current ollama settings from sidebar
                    current_ollama_url = st.session_state.get('ollama_url', 'http://localhost:11434')
                    current_ollama_model = st.session_state.get('ollama_model', 'llama3.1')
                    
                    # Show friendly status
                    if any(word in user_question.lower() for word in ['trend', 'pattern', 'insight', 'analysis']):
                        status_placeholder.info("🔍 Diving deep into your data patterns...")
                    else:
                        status_placeholder.info("🤖 Getting you a friendly answer...")
                    
                    ai_response = chat_with_data(data, user_question, current_ollama_model, current_ollama_url)
                    st.session_state.chat_history.append({"role": "assistant", "content": ai_response})
                    
                    # Clear status
                    status_placeholder.empty()
                    
                except Exception as e:
                    status_placeholder.empty()
                    error_response = f"Oops, I ran into an issue: {str(e)}\n\nTry asking me something like:\n• 'What columns are in this data?'\n• 'What are the key statistics?'\n• 'Tell me about this dataset'"
                    st.session_state.chat_history.append({"role": "assistant", "content": error_response})
        
        # Display chat history
        if st.session_state.chat_history:
            for message in st.session_state.chat_history[-8:]:  # Show last 8 messages
                with st.chat_message(message["role"]):
                    st.write(message["content"])
        else:
            # Show example questions and quick buttons if no chat history
            st.info("""
            💡 **Try asking me questions like:**
            • "What's interesting about my data?"
            • "Tell me about my marketing performance"
            • "Which columns have the highest values?"
            • "How's my data quality?"
            • "What patterns do you see?"
            """)
            
            # Quick action buttons
            st.markdown("**🚀 Quick Questions:**")
            col1, col2, col3 = st.columns(3)
            
            with col1:
                if st.button("📊 Key Statistics", help="Get statistical summary"):
                    st.session_state.chat_history.append({"role": "user", "content": "What are the key statistics?"})
                    st.rerun()
            
            with col2:
                if st.button("📈 Interesting Patterns", help="Find data patterns"):
                    st.session_state.chat_history.append({"role": "user", "content": "What's interesting about my data?"})
                    st.rerun()
                    
            with col3:
                if st.button("📋 Data Overview", help="Show all columns"):
                    st.session_state.chat_history.append({"role": "user", "content": "Give me an overview of this dataset"})
                    st.rerun()
        
        # Clear chat button
        col1, col2 = st.columns([1, 4])
        with col1:
            if st.button("🗑️ Clear Chat"):
                st.session_state.chat_history = []
                st.rerun()
    
    else:
        # Welcome message
        st.info("👆 Load a dataset from the sidebar to get started!")
        
        st.markdown("""
        ### 🎯 New & Improved Features:
        
        #### **🎨 Human-Like Insights**
        - No more corporate jargon or robotic analysis
        - Conversational language that actually makes sense
        - Friendly explanations like you're talking to a colleague
        - Specific observations with real numbers
        
        #### **💬 Smarter Chat**
        - Ask questions in plain English
        - Get helpful, conversational answers
        - Enhanced marketing performance analysis
        - Instant responses with friendly personality
        
        #### **📊 Better Charts**
        - Engaging, descriptive titles
        - Friendly annotations and highlights
        - Up to 8 comprehensive visualizations
        - Performance indicators that make sense
        
        ### 🚀 Getting Started:
        1. **Upload a CSV** or enter a **Kaggle dataset name**
        2. **Generate charts** with human-like insights
        3. **Chat with your data** using natural language
        4. **Download reports** with friendly analysis
        
        ### 💡 Example Questions That Work Great:
        - "What's the story behind my marketing data?"
        - "Which areas need the most attention?"
        - "What patterns should I be excited about?"
        - "Help me understand what's going on here"
        - "What would you focus on if this was your data?"
        
        ### ✨ What Makes This Better:
        - **Conversational tone** instead of business-speak
        - **Specific insights** with actual numbers
        - **Friendly explanations** anyone can understand
        - **Actionable observations** you can actually use
        - **Human personality** that makes data analysis engaging
        """)

if __name__ == "__main__":
    main()