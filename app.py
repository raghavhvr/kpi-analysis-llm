import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import kaggle
import os
from datetime import datetime
import json
import requests
from dotenv import load_dotenv
import base64
import io
from PIL import Image
import tempfile
import warnings

# Suppress all warnings for cleaner output
warnings.filterwarnings('ignore')

# Optional imports that may not be available
try:
    import numpy as np
    import scipy.stats  # Changed from 'from scipy import stats' to avoid conflicts
    HAS_SCIPY = True
except ImportError:
    HAS_SCIPY = False
    st.sidebar.info("💡 Install scipy for statistical trend lines: pip install scipy")

# Load environment variables
load_dotenv()

# Page configuration
st.set_page_config(
    page_title="KPI Insights Dashboard",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Move this function to the top to avoid import issues
def identify_chart_opportunities(df, use_full_dataset=True, sample_size=2000):
    """Identify the best chart types based on data patterns"""
    try:
        numeric_cols = df.select_dtypes(include=['number']).columns.tolist()
        categorical_cols = df.select_dtypes(include=['object', 'category']).columns.tolist()
        datetime_cols = []
        
        # Try to identify datetime columns that are stored as strings (with NO warnings)
        for col in categorical_cols.copy():
            if df[col].dtype == 'object':
                try:
                    # Sample just 5 rows and suppress all warnings
                    sample = df[col].dropna().head(5)
                    if len(sample) > 0:
                        # Try parsing with specific formats first to avoid dateutil warnings
                        with warnings.catch_warnings():
                            warnings.simplefilter("ignore")
                            pd.to_datetime(sample, errors='raise')
                        datetime_cols.append(col)
                        categorical_cols.remove(col)
                except:
                    # Not a datetime column, continue silently
                    pass
        
        opportunities = []
        
        # Handle dataset sampling based on user preference
        if use_full_dataset or len(df) <= sample_size:
            df_sample = df
        else:
            df_sample = df.sample(n=sample_size, random_state=42)
        
        # Time series opportunities (create more if available)
        if datetime_cols and numeric_cols:
            # Create time series for first 2 numeric columns with datetime
            for i, num_col in enumerate(numeric_cols[:2]):
                for j, date_col in enumerate(datetime_cols[:1]):  # Use first datetime column
                    opportunities.append({
                        'type': 'time_series',
                        'title': f'{num_col.replace("_", " ").title()} Over Time',
                        'x': date_col,
                        'y': num_col,
                        'priority': 'high',
                        'data': df_sample
                    })
        
        # KPI metric opportunities - create more distributions
        if numeric_cols:
            count = 0
            for col in numeric_cols[:6]:  # Check first 6 columns instead of 5
                # Look for KPI-like column names
                kpi_keywords = ['revenue', 'sales', 'profit', 'cost', 'price', 'amount', 'value', 
                               'count', 'total', 'avg', 'mean', 'sum', 'rate', 'percent', 'ctr', 'clicks', 'impressions']
                
                if any(keyword in col.lower() for keyword in kpi_keywords) and count < 4:  # Increased to 4
                    opportunities.append({
                        'type': 'kpi_metric',
                        'title': f'{col.replace("_", " ").title()} Distribution',
                        'column': col,
                        'priority': 'high',
                        'data': df_sample
                    })
                    count += 1
        
        # Category performance - create multiple category analyses
        if categorical_cols and numeric_cols:
            chart_count = 0
            for cat_col in categorical_cols[:3]:  # Check first 3 categorical columns
                unique_count = df[cat_col].nunique()
                if 3 <= unique_count <= 15 and chart_count < 3:  # Create up to 3 category charts
                    for num_col in numeric_cols[:2]:  # Use first 2 numeric columns
                        if chart_count < 3:
                            opportunities.append({
                                'type': 'category_performance',
                                'title': f'{num_col.replace("_", " ").title()} by {cat_col.replace("_", " ").title()}',
                                'x': cat_col,
                                'y': num_col,
                                'priority': 'medium',
                                'data': df_sample
                            })
                            chart_count += 1
        
        # Correlation analysis
        if 3 <= len(numeric_cols) <= 8:
            opportunities.append({
                'type': 'correlation',
                'title': 'Key Metrics Correlation Analysis',
                'columns': numeric_cols[:6],  # Up to 6 columns for correlation
                'priority': 'medium',
                'data': df_sample
            })
        
        # Add geographic analysis if detected
        geo_keywords = ['country', 'state', 'city', 'region', 'location', 'address', 'area']
        geo_cols = [col for col in categorical_cols if any(keyword in col.lower() for keyword in geo_keywords)]
        
        if geo_cols and numeric_cols:
            for geo_col in geo_cols[:2]:  # Up to 2 geographic columns
                for num_col in numeric_cols[:2]:  # With first 2 numeric columns
                    if df[geo_col].nunique() <= 20:  # Reasonable number of locations
                        opportunities.append({
                            'type': 'geographic',
                            'title': f'{num_col.replace("_", " ").title()} by {geo_col.replace("_", " ").title()}',
                            'x': geo_col,
                            'y': num_col,
                            'priority': 'medium',
                            'data': df_sample
                        })
        
        # Limit to maximum 8 charts instead of 3
        return opportunities[:8]
        
    except Exception as e:
        st.error(f"Error in identify_chart_opportunities: {str(e)}")
        return []

class SmartChartGenerator:
    """Generate business-relevant charts based on data patterns"""
    
    @staticmethod
    def identify_chart_opportunities(df):
        """Identify the best chart types based on data characteristics"""
        numeric_cols = df.select_dtypes(include=['number']).columns.tolist()
        categorical_cols = df.select_dtypes(include=['object', 'category']).columns.tolist()
        datetime_cols = df.select_dtypes(include=['datetime64']).columns.tolist()
        
        # Try to identify datetime columns that are stored as strings
        for col in categorical_cols.copy():
            if df[col].dtype == 'object':
                try:
                    pd.to_datetime(df[col].dropna().head(100))
                    datetime_cols.append(col)
                    categorical_cols.remove(col)
                except:
                    pass
        
        opportunities = []
        
        # Time series opportunities
        if datetime_cols and numeric_cols:
            for date_col in datetime_cols:
                for num_col in numeric_cols:
                    opportunities.append({
                        'type': 'time_series',
                        'title': f'{num_col} Trend Over Time',
                        'x': date_col,
                        'y': num_col,
                        'priority': 'high'
                    })
        
        # KPI metric opportunities
        if numeric_cols:
            for col in numeric_cols:
                # Look for KPI-like column names
                kpi_keywords = ['revenue', 'sales', 'profit', 'cost', 'price', 'amount', 'value', 
                               'count', 'total', 'avg', 'mean', 'sum', 'rate', 'percent']
                
                if any(keyword in col.lower() for keyword in kpi_keywords):
                    opportunities.append({
                        'type': 'kpi_metric',
                        'title': f'{col} Distribution',
                        'column': col,
                        'priority': 'high'
                    })
        
        # Category performance opportunities
        if categorical_cols and numeric_cols:
            for cat_col in categorical_cols:
                # Limit to reasonable number of categories
                if df[cat_col].nunique() <= 20:
                    for num_col in numeric_cols:
                        opportunities.append({
                            'type': 'category_performance',
                            'title': f'{num_col} by {cat_col}',
                            'x': cat_col,
                            'y': num_col,
                            'priority': 'medium'
                        })
        
        # Correlation opportunities
        if len(numeric_cols) >= 2:
            opportunities.append({
                'type': 'correlation',
                'title': 'Correlation Analysis',
                'columns': numeric_cols[:10],  # Limit to avoid overcrowding
                'priority': 'medium'
            })
        
        # Geographic opportunities
        geo_keywords = ['country', 'state', 'city', 'region', 'location', 'address']
        geo_cols = [col for col in categorical_cols if any(keyword in col.lower() for keyword in geo_keywords)]
        
        if geo_cols and numeric_cols:
            for geo_col in geo_cols:
                for num_col in numeric_cols:
                    opportunities.append({
                        'type': 'geographic',
                        'title': f'{num_col} by {geo_col}',
                        'x': geo_col,
                        'y': num_col,
                        'priority': 'medium'
                    })
        
        return opportunities
    
    @staticmethod
    def create_chart(df, opportunity):
        """Create fast, professional charts with key highlights"""
        try:
            # Use the sampled data from opportunity for faster rendering
            data_to_use = opportunity.get('data', df)
            
            if opportunity['type'] == 'time_series':
                # Fast time series chart with highlights
                df_copy = data_to_use.copy()
                try:
                    with warnings.catch_warnings():
                        warnings.simplefilter("ignore")
                        df_copy[opportunity['x']] = pd.to_datetime(df_copy[opportunity['x']], errors='coerce')
                    df_copy = df_copy.dropna(subset=[opportunity['x']]).sort_values(opportunity['x'])
                    
                    # Create line chart with highlights
                    fig = px.line(df_copy, x=opportunity['x'], y=opportunity['y'],
                                title=opportunity['title'],
                                labels={opportunity['y']: opportunity['y'].replace('_', ' ').title()})
                    
                    # Add high/low point highlights (FAST version)
                    max_idx = df_copy[opportunity['y']].idxmax()
                    min_idx = df_copy[opportunity['y']].idxmin()
                    max_val = df_copy[opportunity['y']].max()
                    min_val = df_copy[opportunity['y']].min()
                    max_date = df_copy.loc[max_idx, opportunity['x']]
                    min_date = df_copy.loc[min_idx, opportunity['x']]
                    
                    # Add annotations for high/low points
                    fig.add_annotation(
                        x=max_date, y=max_val,
                        text=f"📈 High: {max_val:,.0f}",
                        showarrow=True, arrowcolor="green", bgcolor="lightgreen",
                        font=dict(color="darkgreen")
                    )
                    
                    fig.add_annotation(
                        x=min_date, y=min_val,
                        text=f"📉 Low: {min_val:,.0f}",
                        showarrow=True, arrowcolor="red", bgcolor="lightcoral",
                        font=dict(color="darkred")
                    )
                    
                    # Simple trend line if scipy available
                    if HAS_SCIPY and len(df_copy) <= 200:  # Even more restrictive for speed
                        try:
                            x_numeric = range(len(df_copy))
                            slope, intercept, r_value, _, _ = scipy.stats.linregress(x_numeric, df_copy[opportunity['y']])
                            trend_line = [slope * x + intercept for x in x_numeric]
                            
                            fig.add_trace(go.Scatter(
                                x=df_copy[opportunity['x']], y=trend_line,
                                mode='lines', name=f'📊 Trend (R²={r_value**2:.2f})',
                                line=dict(color='orange', width=2, dash='dash')
                            ))
                        except:
                            pass  # Skip trend if error
                    
                    fig.update_layout(template='plotly_white', showlegend=True)
                    return fig
                    
                except Exception:
                    return SmartChartGenerator._create_fallback_chart(data_to_use, opportunity)
                
            elif opportunity['type'] == 'kpi_metric':
                # Fast histogram with mean highlight
                fig = px.histogram(data_to_use, x=opportunity['column'], 
                                 title=opportunity['title'], nbins=15,
                                 labels={opportunity['column']: opportunity['column'].replace('_', ' ').title()})
                
                # Add mean line with value
                mean_val = data_to_use[opportunity['column']].mean()
                median_val = data_to_use[opportunity['column']].median()
                
                fig.add_vline(x=mean_val, line_dash="dash", line_color="red", 
                             annotation_text=f"📊 Avg: {mean_val:.1f}")
                
                fig.update_layout(template='plotly_white', showlegend=False)
                return fig
                
            elif opportunity['type'] == 'category_performance':
                # Fast bar chart with top/bottom highlights
                top_categories = data_to_use[opportunity['x']].value_counts().head(8).index  # Reduced to 8 for speed
                df_filtered = data_to_use[data_to_use[opportunity['x']].isin(top_categories)]
                df_grouped = df_filtered.groupby(opportunity['x'])[opportunity['y']].mean().reset_index()
                df_grouped = df_grouped.sort_values(opportunity['y'], ascending=False)
                
                # Color-code bars (green for best, red for worst)
                colors = []
                for i, _ in enumerate(df_grouped.index):
                    if i == 0:  # Best performer
                        colors.append('lightgreen')
                    elif i == len(df_grouped) - 1:  # Worst performer  
                        colors.append('lightcoral')
                    else:
                        colors.append('lightblue')
                
                fig = px.bar(df_grouped, x=opportunity['x'], y=opportunity['y'],
                           title=opportunity['title'],
                           labels={opportunity['y']: f"Avg {opportunity['y'].replace('_', ' ').title()}"})
                
                # Update colors
                fig.update_traces(marker_color=colors)
                
                # Add value labels on bars
                fig.update_traces(texttemplate='%{y:.1f}', textposition='outside')
                
                # Add annotations for best/worst
                best_val = df_grouped.iloc[0]
                worst_val = df_grouped.iloc[-1]
                
                fig.add_annotation(
                    x=best_val[opportunity['x']], y=best_val[opportunity['y']],
                    text="🏆 Best", showarrow=False, 
                    font=dict(color="darkgreen", size=12)
                )
                
                if len(df_grouped) > 1:
                    fig.add_annotation(
                        x=worst_val[opportunity['x']], y=worst_val[opportunity['y']],
                        text="⚠️ Focus", showarrow=False,
                        font=dict(color="darkred", size=12)
                    )
                
                fig.update_layout(template='plotly_white', xaxis_tickangle=-45, showlegend=False)
                return fig
                
            elif opportunity['type'] == 'correlation':
                # Super fast correlation heatmap
                corr_matrix = data_to_use[opportunity['columns'][:6]].corr()  # Limit to 6 columns max
                
                fig = px.imshow(corr_matrix, 
                              title=opportunity['title'],
                              color_continuous_scale='RdBu_r',
                              aspect="auto")
                
                # Add correlation values as text (simplified)
                fig.update_traces(text=corr_matrix.round(2).values, texttemplate="%{text}")
                fig.update_layout(template='plotly_white')
                return fig
                
            elif opportunity['type'] == 'geographic':
                # Geographic/regional analysis
                top_locations = data_to_use[opportunity['x']].value_counts().head(15).index
                df_filtered = data_to_use[data_to_use[opportunity['x']].isin(top_locations)]
                df_grouped = df_filtered.groupby(opportunity['x'])[opportunity['y']].mean().reset_index()
                df_grouped = df_grouped.sort_values(opportunity['y'], ascending=False)
                
                # Color-code by performance
                colors = []
                for i, _ in enumerate(df_grouped.index):
                    if i == 0:  # Best performer
                        colors.append('#2E8B57')  # Green
                    elif i < 3:  # Top 3
                        colors.append('#4682B4')  # Blue
                    elif i >= len(df_grouped) - 3:  # Bottom 3
                        colors.append('#DC143C')  # Red
                    else:
                        colors.append('#4169E1')  # Royal blue
                
                fig = px.bar(df_grouped, x=opportunity['x'], y=opportunity['y'],
                           title=opportunity['title'],
                           labels={opportunity['y']: f"Avg {opportunity['y'].replace('_', ' ').title()}"})
                
                # Update colors
                fig.update_traces(marker_color=colors)
                
                # Add value labels
                fig.update_traces(texttemplate='%{y:.1f}', textposition='outside')
                
                # Add annotations for best/worst
                best_val = df_grouped.iloc[0]
                fig.add_annotation(
                    x=best_val[opportunity['x']], y=best_val[opportunity['y']],
                    text="🏆", showarrow=False, font=dict(size=16)
                )
                
                if len(df_grouped) > 1:
                    worst_val = df_grouped.iloc[-1]
                    fig.add_annotation(
                        x=worst_val[opportunity['x']], y=worst_val[opportunity['y']],
                        text="⚠️", showarrow=False, font=dict(size=16)
                    )
                
                fig.update_layout(template='plotly_white', xaxis_tickangle=-45, showlegend=False)
                return fig
                
        except Exception as e:
            st.warning(f"⚠️ Chart generation issue (using fallback): {str(e)}")
            return SmartChartGenerator._create_fallback_chart(data_to_use, opportunity)
        
        return None
    
    @staticmethod
    def _create_fallback_chart(df, opportunity):
        """Super simple fallback charts"""
        try:
            if opportunity['type'] == 'time_series':
                fig = px.scatter(df.sample(min(100, len(df))), 
                               x=opportunity['x'], y=opportunity['y'],
                               title=f"Sample: {opportunity['title']}")
                return fig
                
            elif opportunity['type'] == 'kpi_metric':
                fig = px.histogram(df.sample(min(100, len(df))), 
                                 x=opportunity['column'], title=f"Sample: {opportunity['title']}")
                return fig
                
            elif opportunity['type'] == 'category_performance':
                sample_df = df.sample(min(100, len(df)))
                fig = px.bar(sample_df.value_counts(opportunity['x']).head(5).reset_index(),
                           x='index', y=opportunity['x'], title=f"Top Categories: {opportunity['title']}")
                return fig
                
        except Exception:
            # Ultimate fallback - simple bar chart of value counts
            try:
                numeric_cols = df.select_dtypes(include=['number']).columns
                if len(numeric_cols) > 0:
                    col = numeric_cols[0]
                    fig = px.histogram(df.sample(min(50, len(df))), x=col, title="Data Sample")
                    return fig
            except:
                pass
        
        return None

class LLMInsights:
    """Generate AI-powered insights using Ollama (open source)"""
    
    @staticmethod
    def analyze_data_with_llm(df, chart_info, ollama_model='llama3.1', ollama_url='http://localhost:11434'):
        """Generate insights using Ollama with timeout"""
        
        # Create meaningful data context for LLM
        numeric_cols = df.select_dtypes(include=['number']).columns.tolist()
        categorical_cols = df.select_dtypes(include=['object']).columns.tolist()
        
        # Get actual data patterns for better insights
        insights_context = LLMInsights._build_data_context(df, chart_info)
        
        prompt = f"""You are a senior business analyst reviewing data for a stakeholder presentation. 

{insights_context}

Chart: {chart_info.get('title', 'Data Analysis')}

Write 2-3 clear, actionable business insights in natural language that would be valuable for decision makers. Focus on:
- What the data tells us about business performance
- Key trends or patterns that matter
- Specific recommendations or areas for investigation

Write as if explaining to a business executive. Be concise but insightful. Avoid technical jargon.

Example format:
"Revenue shows strong growth in Q2, with the North region outperforming by 23%. However, customer satisfaction dipped slightly in June, suggesting we should investigate service quality during our busiest period."
"""
        
        try:
            # Try Ollama API with short timeout
            response = requests.post(
                f"{ollama_url}/api/generate",
                json={
                    'model': ollama_model,
                    'prompt': prompt,
                    'stream': False,
                    'options': {
                        'temperature': 0.7,
                        'top_p': 0.9,
                        'num_predict': 150  # Reduced for faster response
                    }
                },
                timeout=8  # Very short timeout for insights
            )
            
            if response.status_code == 200:
                result = response.json()
                return result.get('response', '').strip()
            
        except (requests.exceptions.RequestException, requests.exceptions.Timeout):
            # Quick fallback - don't wait for Ollama
            pass
        except Exception as e:
            # Any other error, use fallback
            pass
        
        # Enhanced fallback insights (always available)
        return LLMInsights.generate_business_insights(df, chart_info)
    
    @staticmethod
    def _build_data_context(df, chart_info):
        """Build meaningful context about the data for LLM"""
        context_parts = []
        
        # Dataset basics
        context_parts.append(f"Dataset: {df.shape[0]:,} records with {df.shape[1]} variables")
        
        # Key metrics summary
        numeric_cols = df.select_dtypes(include=['number']).columns
        if len(numeric_cols) > 0:
            key_stats = []
            for col in numeric_cols[:3]:  # Top 3 numeric columns
                mean_val = df[col].mean()
                median_val = df[col].median()
                if mean_val > 1000:
                    key_stats.append(f"{col}: avg ${mean_val:,.0f}")
                elif mean_val > 1:
                    key_stats.append(f"{col}: avg {mean_val:.1f}")
                else:
                    key_stats.append(f"{col}: avg {mean_val:.3f}")
            
            if key_stats:
                context_parts.append("Key metrics: " + ", ".join(key_stats))
        
        # Category insights
        categorical_cols = df.select_dtypes(include=['object']).columns
        if len(categorical_cols) > 0:
            cat_insights = []
            for col in categorical_cols[:2]:
                top_cat = df[col].value_counts().index[0]
                top_pct = (df[col].value_counts().iloc[0] / len(df)) * 100
                cat_insights.append(f"{col} dominated by {top_cat} ({top_pct:.0f}%)")
            
            if cat_insights:
                context_parts.append(" • ".join(cat_insights))
        
        return "\n".join(context_parts)
    
    @staticmethod
    def generate_business_insights(df, chart_info):
        """Generate human, business-focused insights without LLM"""
        insights = []
        
        chart_title = chart_info.get('title', '').lower()
        numeric_cols = df.select_dtypes(include=['number']).columns
        categorical_cols = df.select_dtypes(include=['object']).columns
        
        # Time series insights
        if 'trend' in chart_title or 'time' in chart_title:
            if len(numeric_cols) > 0:
                col = numeric_cols[0]
                try:
                    # Simple trend analysis
                    first_half = df[col][:len(df)//2].mean()
                    second_half = df[col][len(df)//2:].mean()
                    change_pct = ((second_half - first_half) / first_half) * 100
                    
                    if abs(change_pct) > 10:
                        direction = "increased" if change_pct > 0 else "decreased"
                        insights.append(f"The trend shows {col.replace('_', ' ')} has {direction} by {abs(change_pct):.1f}% over the time period, indicating {'strong growth' if change_pct > 0 else 'concerning decline'}.")
                    else:
                        insights.append(f"{col.replace('_', ' ').title()} remains relatively stable over time, suggesting consistent performance without major volatility.")
                except:
                    insights.append(f"The {chart_title} reveals patterns over time that warrant executive attention for strategic planning.")
        
        # Performance by category insights
        elif 'by' in chart_title and len(categorical_cols) > 0 and len(numeric_cols) > 0:
            cat_col = categorical_cols[0]
            num_col = numeric_cols[0]
            
            try:
                performance = df.groupby(cat_col)[num_col].mean().sort_values(ascending=False)
                top_performer = performance.index[0]
                bottom_performer = performance.index[-1]
                performance_gap = ((performance.iloc[0] - performance.iloc[-1]) / performance.iloc[-1]) * 100
                
                insights.append(f"{top_performer} leads in {num_col.replace('_', ' ')} performance, outperforming {bottom_performer} by {performance_gap:.0f}%. This suggests successful strategies in {top_performer} that could be replicated elsewhere.")
                
                # Check for consistency
                cv = performance.std() / performance.mean()
                if cv > 0.3:
                    insights.append(f"There's significant variation in {num_col.replace('_', ' ')} across different {cat_col.replace('_', ' ')} categories, indicating potential optimization opportunities.")
                
            except:
                insights.append(f"Performance varies significantly across {cat_col.replace('_', ' ')} categories, suggesting targeted strategies may be needed for different segments.")
        
        # Distribution insights
        elif 'distribution' in chart_title:
            if len(numeric_cols) > 0:
                col = numeric_cols[0]
                try:
                    q75 = df[col].quantile(0.75)
                    q25 = df[col].quantile(0.25)
                    median = df[col].median()
                    mean = df[col].mean()
                    
                    # Skewness check
                    if mean > median * 1.2:
                        insights.append(f"The {col.replace('_', ' ')} distribution is right-skewed, with a few high performers driving up the average. Focus on understanding what drives these top performers.")
                    elif mean < median * 0.8:
                        insights.append(f"The {col.replace('_', ' ')} distribution shows most values are above average, but some underperformers are pulling down the mean. Investigate bottom quartile for improvement opportunities.")
                    else:
                        insights.append(f"{col.replace('_', ' ').title()} shows a fairly balanced distribution, indicating consistent performance across the dataset.")
                        
                except:
                    insights.append(f"The distribution of {col.replace('_', ' ')} provides insights into performance patterns that could inform resource allocation decisions.")
        
        # Correlation insights
        elif 'correlation' in chart_title:
            if len(numeric_cols) >= 2:
                try:
                    corr_matrix = df[numeric_cols].corr()
                    # Find strongest correlations (excluding self-correlation)
                    strong_corrs = []
                    for i in range(len(corr_matrix.columns)):
                        for j in range(i+1, len(corr_matrix.columns)):
                            corr_val = corr_matrix.iloc[i, j]
                            if abs(corr_val) > 0.7:
                                col1, col2 = corr_matrix.columns[i], corr_matrix.columns[j]
                                relationship = "strongly positive" if corr_val > 0 else "strongly negative"
                                strong_corrs.append(f"{col1.replace('_', ' ')} and {col2.replace('_', ' ')} show a {relationship} relationship")
                    
                    if strong_corrs:
                        insights.append(f"Key finding: {strong_corrs[0]}. This suggests these metrics move together and could be leveraged for predictive planning.")
                    else:
                        insights.append("The correlation analysis shows moderate relationships between variables, suggesting each metric provides unique insights for decision-making.")
                        
                except:
                    insights.append("The correlation patterns reveal important relationships between key business metrics that should inform strategic planning.")
        
        # Default business insight
        if len(insights) == 0:
            if len(df) > 1000:
                insights.append(f"This comprehensive dataset of {len(df):,} records provides robust insights for data-driven decision making.")
            else:
                insights.append(f"This focused dataset offers targeted insights that can inform immediate business actions.")
            
            # Data quality insight
            missing_pct = (df.isnull().sum().sum() / (df.shape[0] * df.shape[1])) * 100
            if missing_pct < 5:
                insights.append("Data quality is excellent, providing confidence in the analytical findings and recommendations.")
            elif missing_pct < 15:
                insights.append("Data quality is good overall, though some targeted data collection improvements could enhance future analysis.")
        
        # Ensure we have actionable insights
        if len(insights) < 2:
            insights.append("These findings warrant further investigation to identify specific action items and optimization opportunities.")
        
        return " ".join(insights[:3])  # Limit to 3 insights for readability

class KPIDashboard:
    def __init__(self):
        self.data = None
        self.charts = {}
        self.data_source = None
        
    def reset_data(self):
        """Reset all data and charts"""
        self.data = None
        self.charts = {}
        self.data_source = None
        
    def load_csv_upload(self, uploaded_file):
        """Load CSV from uploaded file"""
        try:
            self.data = pd.read_csv(uploaded_file)
            self.data_source = uploaded_file.name
            return True, f"Successfully loaded {uploaded_file.name}"
        except Exception as e:
            return False, f"Error loading CSV: {str(e)}"
        
    def load_kaggle_dataset(self, dataset_name):
        """Load dataset from Kaggle"""
        try:
            # Reset previous data
            self.reset_data()
            
            # Download dataset
            kaggle.api.dataset_download_files(dataset_name, path='./data', unzip=True)
            
            # Find CSV files in the data directory
            data_files = [f for f in os.listdir('./data') if f.endswith('.csv')]
            
            if data_files:
                # Load the first CSV file found
                file_path = f"./data/{data_files[0]}"
                self.data = pd.read_csv(file_path)
                self.data_source = f"Kaggle: {dataset_name}"
                return True, f"Successfully loaded {data_files[0]} from Kaggle"
            else:
                return False, "No CSV files found in the dataset"
                
        except Exception as e:
            return False, f"Error loading dataset: {str(e)}"
    
    def generate_smart_charts(self, use_full_dataset=None, sample_size=None):
        """Generate intelligent, business-relevant charts with user preferences"""
        if self.data is None:
            return {}
        
        # Use defaults if not provided
        if use_full_dataset is None:
            use_full_dataset = len(self.data) <= 5000
        if sample_size is None:
            sample_size = 2000
        
        # Reset charts
        self.charts = {}
        
        # Get chart opportunities with user preferences using the global function
        opportunities = identify_chart_opportunities(
            self.data, use_full_dataset, sample_size
        )
        
        if not opportunities:
            st.warning("No suitable chart opportunities found in this dataset.")
            return {}
        
        # Create charts without progress bars for maximum speed
        for i, opportunity in enumerate(opportunities):
            try:
                chart = SmartChartGenerator.create_chart(self.data, opportunity)
                if chart:
                    self.charts[f'chart_{i+1}'] = {
                        'title': opportunity['title'],
                        'figure': chart,
                        'opportunity': opportunity
                    }
            except Exception as e:
                st.error(f"❌ Error creating {opportunity['title']}: {str(e)}")
        
        return self.charts
    
    def create_powerpoint_with_charts(self, charts, insights):
        """Create PowerPoint with actual chart images"""
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "KPI Analysis Report"
        subtitle.text = f"Data Source: {self.data_source}\nGenerated: {datetime.now().strftime('%B %d, %Y')}"
        
        # Executive Summary slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Executive Summary"
        tf = content.text_frame
        tf.clear()
        
        # Add data overview
        p = tf.paragraphs[0]
        p.text = f"Dataset Overview: {self.data.shape[0]:,} rows, {self.data.shape[1]} columns"
        
        # Add key insights
        for insight in insights[:4]:
            p = tf.add_paragraph()
            p.text = f"• {insight['insight'][:100]}..."
            p.level = 0
        
        # Individual chart slides
        for chart_name, chart_info in charts.items():
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
            
            # Add title
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_shape.text_frame
            title_frame.text = chart_info['title']
            title_frame.paragraphs[0].font.size = Inches(0.3)
            title_frame.paragraphs[0].font.bold = True
            
            # Save chart as image and add to slide
            try:
                # Save plotly chart as image
                img_bytes = chart_info['figure'].to_image(format="png", width=800, height=500)
                
                # Create temporary file
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                    tmp_file.write(img_bytes)
                    tmp_path = tmp_file.name
                
                # Add image to slide
                slide.shapes.add_picture(tmp_path, Inches(1), Inches(1.5), Inches(8), Inches(5))
                
                # Clean up temp file
                os.unlink(tmp_path)
                
                # Add insight text
                insight_text = next((insight['insight'] for insight in insights if insight['chart'] == chart_name), 
                                  "Key business insights can be derived from this visualization.")
                
                text_box = slide.shapes.add_textbox(Inches(1), Inches(7), Inches(8), Inches(1))
                text_frame = text_box.text_frame
                text_frame.text = f"Insight: {insight_text}"
                text_frame.word_wrap = True
                
            except Exception as e:
                # Fallback: add text-only slide
                content_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
                content_frame = content_shape.text_frame
                content_frame.text = f"Chart: {chart_info['title']}\n\nNote: Chart image could not be embedded.\nError: {str(e)}"
        
        # Save presentation
        os.makedirs('generated_ppts', exist_ok=True)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"generated_ppts/KPI_Report_{timestamp}.pptx"
        prs.save(filename)
        
        return filename

def get_basic_data_info(data, question):
    """Provide basic data information without AI"""
    question_lower = question.lower()
    
    # Basic data info
    numeric_cols = data.select_dtypes(include=['number']).columns.tolist()
    categorical_cols = data.select_dtypes(include=['object']).columns.tolist()
    
    # Marketing-specific analysis
    if any(word in question_lower for word in ['marketing', 'campaign', 'ctr', 'conversion', 'performance']):
        marketing_insights = []
        
        # CTR analysis
        ctr_cols = [col for col in data.columns if 'ctr' in col.lower()]
        if ctr_cols:
            ctr_col = ctr_cols[0]
            avg_ctr = data[ctr_col].mean()
            max_ctr = data[ctr_col].max()
            min_ctr = data[ctr_col].min()
            marketing_insights.append(f"📊 **{ctr_col}**: Avg {avg_ctr:.2f}%, Range {min_ctr:.2f}% to {max_ctr:.2f}%")
        
        # Impressions analysis
        impression_cols = [col for col in data.columns if 'impression' in col.lower()]
        if impression_cols:
            imp_col = impression_cols[0]
            avg_imp = data[imp_col].mean()
            max_imp = data[imp_col].max()
            marketing_insights.append(f"👁️ **{imp_col}**: Avg {avg_imp:,.0f}, Peak {max_imp:,.0f}")
        
        # Click analysis
        click_cols = [col for col in data.columns if 'click' in col.lower()]
        if click_cols:
            click_col = click_cols[0]
            avg_clicks = data[click_col].mean()
            total_clicks = data[click_col].sum()
            marketing_insights.append(f"🖱️ **{click_col}**: Avg {avg_clicks:,.0f}, Total {total_clicks:,.0f}")
        
        # Campaign type analysis
        campaign_cols = [col for col in categorical_cols if any(word in col.lower() for word in ['campaign', 'type', 'channel', 'source'])]
        if campaign_cols:
            camp_col = campaign_cols[0]
            top_campaigns = data[camp_col].value_counts().head(3)
            campaign_list = ', '.join([f"{idx} ({count})" for idx, count in top_campaigns.items()])
            marketing_insights.append(f"🎯 **{camp_col}**: Top performers - {campaign_list}")
        
        if marketing_insights:
            return f"""🎯 **Marketing Performance Overview:**

{chr(10).join(marketing_insights)}

📈 **Quick Insights:**
• Dataset contains {data.shape[0]:,} marketing records
• {len([col for col in numeric_cols if any(kw in col.lower() for kw in ['ctr', 'click', 'impression', 'conversion'])])} key marketing metrics detected
• For detailed trend analysis, generate charts above!"""
    
    if any(word in question_lower for word in ['columns', 'column', 'fields']):
        return f"""📊 **Dataset Columns:**

**Numeric columns ({len(numeric_cols)}):** {', '.join(numeric_cols[:10])}{'...' if len(numeric_cols) > 10 else ''}

**Category columns ({len(categorical_cols)}):** {', '.join(categorical_cols[:10])}{'...' if len(categorical_cols) > 10 else ''}

**Total:** {len(data.columns)} columns"""

    elif any(word in question_lower for word in ['rows', 'size', 'shape', 'big', 'large']):
        missing = data.isnull().sum().sum()
        complete = ((data.shape[0] * data.shape[1]) - missing) / (data.shape[0] * data.shape[1]) * 100
        return f"""📏 **Dataset Size:**

• **Rows:** {data.shape[0]:,}
• **Columns:** {data.shape[1]}
• **Total cells:** {data.shape[0] * data.shape[1]:,}
• **Missing values:** {missing:,}
• **Data completeness:** {complete:.1f}%"""

    elif any(word in question_lower for word in ['stats', 'statistics', 'summary', 'describe']):
        if numeric_cols:
            first_col = numeric_cols[0]
            stats = data[first_col].describe()
            return f"""📈 **Key Statistics for {first_col}:**

• **Average:** {stats['mean']:.2f}
• **Median:** {stats['50%']:.2f}
• **Min:** {stats['min']:.2f}
• **Max:** {stats['max']:.2f}
• **Standard deviation:** {stats['std']:.2f}

*For detailed analysis of all columns, generate charts using the button above.*"""

    elif any(word in question_lower for word in ['missing', 'null', 'empty', 'quality']):
        missing_by_col = data.isnull().sum()
        problematic = missing_by_col[missing_by_col > 0].head(5)
        if len(problematic) > 0:
            missing_info = '\n'.join([f"• {col}: {count:,} missing ({count/len(data)*100:.1f}%)" for col, count in problematic.items()])
            return f"""🔍 **Data Quality Check:**

**Columns with missing values:**
{missing_info}

**Overall data completeness:** {((data.shape[0] * data.shape[1]) - data.isnull().sum().sum()) / (data.shape[0] * data.shape[1]) * 100:.1f}%"""
        else:
            return "✅ **Excellent data quality!** No missing values found in any columns."

    elif any(word in question_lower for word in ['trend', 'pattern', 'insight', 'analysis']):
        # Enhanced trend analysis for marketing data
        trend_insights = []
        
        # Analyze numeric trends
        if numeric_cols:
            for col in numeric_cols[:3]:
                series = data[col].dropna()
                if len(series) > 10:
                    # Simple trend analysis
                    first_quarter = series.iloc[:len(series)//4].mean()
                    last_quarter = series.iloc[-len(series)//4:].mean()
                    change = ((last_quarter - first_quarter) / first_quarter) * 100 if first_quarter > 0 else 0
                    
                    if abs(change) > 10:
                        direction = "increasing" if change > 0 else "decreasing"
                        trend_insights.append(f"📈 **{col}**: {direction} trend with {abs(change):.1f}% change")
                    else:
                        trend_insights.append(f"📊 **{col}**: Stable pattern (±{abs(change):.1f}%)")
        
        # Analyze categorical patterns
        if categorical_cols:
            for col in categorical_cols[:2]:
                if data[col].nunique() <= 20:
                    top_cat = data[col].value_counts()
                    dominant = top_cat.iloc[0]
                    total = len(data)
                    dominance_pct = (dominant / total) * 100
                    trend_insights.append(f"🏷️ **{col}**: {top_cat.index[0]} dominates ({dominance_pct:.1f}% of data)")
        
        # Marketing-specific insights
        marketing_keywords = ['ctr', 'click', 'impression', 'conversion', 'spend', 'cost']
        marketing_cols = [col for col in numeric_cols if any(kw in col.lower() for kw in marketing_keywords)]
        
        if marketing_cols:
            trend_insights.append(f"🎯 **Marketing Focus**: Key metrics detected - {', '.join(marketing_cols[:3])}")
            
            # CTR analysis if available
            ctr_cols = [col for col in data.columns if 'ctr' in col.lower()]
            if ctr_cols:
                ctr_col = ctr_cols[0]
                avg_ctr = data[ctr_col].mean()
                high_ctr = data[ctr_col].quantile(0.9)
                trend_insights.append(f"📊 **CTR Performance**: Average {avg_ctr:.2f}%, top 10% achieve {high_ctr:.2f}%")
        
        if trend_insights:
            insights_text = '\n'.join(trend_insights[:5])  # Limit to top 5 insights
            return f"""🎯 **Key Trends Identified:**

{insights_text}

📊 **To get detailed trend charts:**
1. Use the "Generate Charts" button above
2. Look for time series and correlation visualizations
3. High/low points will be automatically marked

💡 **For deeper insights:** Set up Ollama for AI-powered trend analysis!"""
        else:
            return f"""🎯 **Trend Analysis Summary:**

Your dataset contains {data.shape[0]:,} records with {len(numeric_cols)} numeric and {len(categorical_cols)} categorical variables. 

📊 **To identify trends:**
1. Use the "Generate Comprehensive Charts" button above
2. Time series charts will show patterns over time
3. Category charts will reveal performance differences
4. Correlation maps will show relationships between variables

💡 **Try asking more specific questions like:**
• "What's the average CTR?"
• "Which column has the highest values?"
• "Show me data quality info"""

    else:
        return f"""🤖 **I can help you explore your data!**

**Dataset Overview:** {data.shape[0]:,} rows × {data.shape[1]} columns

**Available questions:**
• "What columns are in this data?"
• "How big is this dataset?"
• "Show me key statistics"
• "Check data quality"
• "What insights can you provide?"

**For advanced analysis:** Generate charts above or set up Ollama for AI-powered insights!"""
def chat_with_data(data, question, ollama_model='llama3.1', ollama_url='http://localhost:11434'):
    """Chat interface for data exploration using Ollama with fallback"""
    if data is None:
        return "Please load a dataset first to start chatting with your data."
    
    # Create comprehensive data context
    numeric_cols = data.select_dtypes(include=['number']).columns.tolist()
    categorical_cols = data.select_dtypes(include=['object']).columns.tolist()
    
    # Sample of actual data for context (smaller sample for speed)
    data_sample = data.head(2).to_string() if len(data) > 0 else "No data available"
    
    # Statistical summaries for numeric columns (limit for speed)
    numeric_summary = ""
    if numeric_cols:
        for col in numeric_cols[:3]:
            mean_val = data[col].mean()
            max_val = data[col].max()
            min_val = data[col].min()
            numeric_summary += f"• {col}: range {min_val:.1f} to {max_val:.1f}, average {mean_val:.1f}\n"
    
    # Category summaries (limit for speed)
    category_summary = ""
    if categorical_cols:
        for col in categorical_cols[:2]:
            top_categories = data[col].value_counts().head(2)
            category_summary += f"• {col}: top values are {', '.join([f'{idx} ({count})' for idx, count in top_categories.items()])}\n"
    
    data_context = f"""You are analyzing a business dataset with {data.shape[0]} rows and {data.shape[1]} columns.

Dataset Structure:
- Numeric columns: {', '.join(numeric_cols[:5]) if numeric_cols else 'None'}
- Category columns: {', '.join(categorical_cols[:5]) if categorical_cols else 'None'}

Key Statistics:
{numeric_summary[:200]}

Top Categories:
{category_summary[:200]}

User Question: {question}

Provide a helpful, business-focused answer in 2-3 sentences. Be conversational and practical."""
    
    try:
        # Try Ollama API with longer timeout for analytical questions
        timeout_duration = 15 if any(word in question.lower() for word in ['trend', 'pattern', 'insight', 'analysis', 'correlation']) else 8
        
        response = requests.post(
            f"{ollama_url}/api/generate",
            json={
                'model': ollama_model,
                'prompt': data_context,
                'stream': False,
                'options': {
                    'temperature': 0.3,
                    'top_p': 0.9,
                    'num_predict': 150
                }
            },
            timeout=timeout_duration  # Longer timeout for analytical questions
        )
        
        if response.status_code == 200:
            result = response.json()
            ai_response = result.get('response', '').strip()
            if ai_response and len(ai_response) > 10:  # Valid response
                return ai_response
            else:
                # Empty response, use fallback
                return get_basic_data_info(data, question)
        else:
            # Bad status code, use fallback
            return get_basic_data_info(data, question)
            
    except requests.exceptions.ConnectionError:
        # Ollama not running - provide helpful info
        basic_info = get_basic_data_info(data, question)
        return f"""{basic_info}

🤖 **Want AI-powered insights?**
1. Install Ollama: https://ollama.ai
2. Run: `ollama serve`
3. Download model: `ollama pull llama3.1`"""
        
    except requests.exceptions.Timeout:
        # Timeout - provide enhanced basic response for analytical questions
        if any(word in question.lower() for word in ['trend', 'pattern', 'insight', 'analysis']):
            return get_basic_data_info(data, question)  # This now has enhanced trend analysis
        else:
            return f"""⏰ AI response timed out. Here's what I can tell you:

{get_basic_data_info(data, question)}

**Tip:** For instant responses, try questions like "What columns are available?" or "Show me key statistics"."""
        
    except Exception as e:
        # Any other error - provide basic analysis
        return get_basic_data_info(data, question)

def main():
    st.title("📊 KPI Insights Dashboard v2.4")
    st.markdown("**Comprehensive Analysis Version** - Generate up to 8 charts with working AI chat")
    
    # Initialize dashboard
    if 'dashboard' not in st.session_state:
        st.session_state.dashboard = KPIDashboard()
    
    # Sidebar
    st.sidebar.title("🚀 Dashboard Controls")
    
    # Step 1: Data Loading
    st.sidebar.header("1. Load Data")
    
    # Data source selection
    data_source = st.sidebar.radio(
        "Choose data source:",
        ["Upload CSV", "Kaggle Dataset"]
    )
    
    if data_source == "Upload CSV":
        uploaded_file = st.sidebar.file_uploader("Choose a CSV file", type="csv")
        if uploaded_file is not None:
            if st.sidebar.button("Load CSV"):
                with st.spinner("Loading CSV file..."):
                    success, message = st.session_state.dashboard.load_csv_upload(uploaded_file)
                    if success:
                        st.sidebar.success(message)
                        st.rerun()
                    else:
                        st.sidebar.error(message)
    
    else:  # Kaggle Dataset
        dataset_name = st.sidebar.text_input(
            "Enter Kaggle Dataset Name",
            placeholder="e.g., username/dataset-name",
            help="Format: username/dataset-name (as shown in Kaggle URL)"
        )
        
        if st.sidebar.button("Load Kaggle Dataset"):
            if dataset_name:
                with st.spinner("Loading dataset from Kaggle..."):
                    success, message = st.session_state.dashboard.load_kaggle_dataset(dataset_name)
                    if success:
                        st.sidebar.success(message)
                        st.rerun()
                    else:
                        st.sidebar.error(message)
            else:
                st.sidebar.error("Please enter a dataset name")
        
        # Sample datasets for testing
        st.sidebar.subheader("Sample Datasets:")
        sample_datasets = [
            "russellyates88/stock-market-data",
            "prasadperera/the-boston-housing-dataset",
            "vikramtiwari/pima-indians-diabetes-database"
        ]
        
        for dataset in sample_datasets:
            if st.sidebar.button(f"📊 {dataset.split('/')[1][:20]}...", key=dataset):
                with st.spinner(f"Loading {dataset}..."):
                    success, message = st.session_state.dashboard.load_kaggle_dataset(dataset)
                    if success:
                        st.sidebar.success(message)
                        st.rerun()
                    else:
                        st.sidebar.error(message)
    
    # Ollama Configuration
    st.sidebar.header("2. AI Configuration")
    
    with st.sidebar.expander("🤖 Ollama Setup", expanded=False):
        st.markdown("""
        **For Enhanced AI Insights:**
        1. Install Ollama: https://ollama.ai
        2. Run: `ollama serve`
        3. Download model: `ollama pull llama3.1`
        """)
    
    ollama_url = st.sidebar.text_input("Ollama URL", value="http://localhost:11434", 
                                      help="Default: http://localhost:11434")
    st.session_state.ollama_url = ollama_url
    
    ollama_model = st.sidebar.selectbox("Ollama Model", 
                                       ["llama3.1", "llama3", "mistral", "codellama"],
                                       help="Make sure the model is downloaded: ollama pull <model>")
    st.session_state.ollama_model = ollama_model
    
    # Analysis Options
    st.sidebar.header("3. Analysis Options")
    
    # Set initial values
    if st.session_state.dashboard.data is not None:
        default_full_dataset = len(st.session_state.dashboard.data) <= 5000
        max_sample = min(5000, len(st.session_state.dashboard.data))
    else:
        default_full_dataset = True
        max_sample = 5000
    
    use_full_dataset = st.sidebar.checkbox(
        "Use Full Dataset", 
        value=default_full_dataset,
        help="Uncheck for faster analysis with sampling"
    )
    st.session_state.use_full_dataset = use_full_dataset
    
    if not use_full_dataset and st.session_state.dashboard.data is not None:
        sample_size = st.sidebar.slider(
            "Sample Size", 
            min_value=500, 
            max_value=max_sample, 
            value=min(2000, max_sample),
            help="Smaller = Faster, Larger = More accurate"
        )
        st.session_state.sample_size = sample_size
        st.sidebar.info(f"Will analyze {sample_size:,} out of {len(st.session_state.dashboard.data):,} rows")
    else:
        st.session_state.sample_size = 2000  # Default when using full dataset
    
    # Test Ollama connection
    if st.sidebar.button("🔌 Test Ollama Connection"):
        try:
            response = requests.get(f"{ollama_url}/api/tags", timeout=5)
            if response.status_code == 200:
                models = response.json().get('models', [])
                if models:
                    available_models = [m['name'] for m in models]
                    st.sidebar.success(f"✅ Connected! Available models: {', '.join(available_models[:3])}")
                else:
                    st.sidebar.warning("Ollama connected but no models found. Run: ollama pull llama3.1")
            else:
                st.sidebar.error("❌ Ollama not responding")
        except:
            st.sidebar.error("❌ Cannot connect to Ollama. Make sure it's running!")
            st.sidebar.info("Start with: `ollama serve`")
    
    # Main content area
    if st.session_state.dashboard.data is not None:
        data = st.session_state.dashboard.data
        
        # Data preview
        st.header("📋 Data Overview")
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            st.metric("📊 Total Rows", f"{len(data):,}")
        with col2:
            st.metric("📈 Total Columns", len(data.columns))
        with col3:
            st.metric("❌ Missing Values", f"{data.isnull().sum().sum():,}")
        with col4:
            st.metric("💾 Memory Usage", f"{data.memory_usage(deep=True).sum() / 1024**2:.1f} MB")
        
        # Data source info
        st.info(f"**Data Source:** {st.session_state.dashboard.data_source}")
        
        # Show data sample
        with st.expander("📊 View Data Sample", expanded=False):
            st.dataframe(data.head(100), width='stretch')
        
        # Generate charts with better feedback
        if st.button("🚀 Generate Comprehensive Charts & Analysis", type="primary"):
            # Quick data validation
            if len(data) == 0:
                st.error("❌ Dataset is empty! Please load a different dataset.")
                return
            
            if len(data.columns) < 2:
                st.error("❌ Dataset needs at least 2 columns for meaningful analysis.")
                return
            
            # Show what we're working with
            numeric_cols = data.select_dtypes(include=['number']).columns.tolist()
            categorical_cols = data.select_dtypes(include=['object']).columns.tolist()
            
            with st.spinner("🎨 Creating charts..."):
                # Get user preferences from sidebar (with fallbacks)
                try:
                    use_full = st.session_state.get('use_full_dataset', len(data) <= 5000)
                    sample_sz = st.session_state.get('sample_size', 2000)
                except:
                    use_full = len(data) <= 5000
                    sample_sz = 2000
                
                # Show dataset info based on user preferences
                if use_full or len(data) <= 5000:
                    analysis_scope = "Full dataset"
                    analysis_size = len(data)
                else:
                    analysis_scope = f"Sample of {sample_sz:,} rows"
                    analysis_size = min(sample_sz, len(data))
                
                st.info(f"""
                ⚡ **Chart Generation:**
                - Dataset size: {len(data):,} rows
                - Analysis scope: {analysis_scope}
                - Processing: {analysis_size:,} rows for visualization
                - Charts: Up to 8 comprehensive visualizations with highlights
                - Expected time: 10-15 seconds
                """)
                
                charts = st.session_state.dashboard.generate_smart_charts(
                    use_full_dataset=use_full,
                    sample_size=sample_sz
                )
            
            if charts:
                # Add option to skip AI insights for faster generation
                col1, col2 = st.columns([3, 1])
                with col1:
                    st.success(f"✅ Generated {len(charts)} charts with highlights!")
                with col2:
                    skip_insights = st.button("⚡ Skip AI Insights", help="Show charts without AI analysis for speed")
                
                # Generate insights based on user choice
                if skip_insights:
                    # Create simple insights without AI
                    insights = []
                    for chart_name, chart_info in charts.items():
                        insights.append({
                            'chart': chart_name,
                            'title': chart_info['title'],
                            'insight': f"This {chart_info['title'].lower()} provides valuable business insights for strategic decision-making. The visualization highlights key patterns and performance indicators that warrant further investigation."
                        })
                    st.info("⚡ Charts ready! AI insights skipped for maximum speed.")
                else:
                    # Generate AI insights with timeout
                    with st.spinner("🤖 Generating AI insights (max 15 seconds)..."):
                        insights = []
                        insight_timeout = 15  # 15 seconds total for all insights
                        start_time = datetime.now()
                        
                        for chart_name, chart_info in charts.items():
                            if (datetime.now() - start_time).seconds > insight_timeout:
                                # Timeout reached, use simple insights for remaining charts
                                insights.append({
                                    'chart': chart_name,
                                    'title': chart_info['title'],
                                    'insight': f"Analysis of {chart_info['title'].lower()} shows important patterns that can inform business strategy and operational decisions."
                                })
                                continue
                                
                            try:
                                insight_text = LLMInsights.analyze_data_with_llm(
                                    data, chart_info, ollama_model, ollama_url
                                )
                                insights.append({
                                    'chart': chart_name,
                                    'title': chart_info['title'],
                                    'insight': insight_text
                                })
                            except Exception:
                                # Fallback insight
                                insights.append({
                                    'chart': chart_name,
                                    'title': chart_info['title'],
                                    'insight': f"The {chart_info['title'].lower()} reveals key business patterns and performance indicators that provide valuable insights for decision-making."
                                })
                
                # Store in session state
                st.session_state.charts = charts
                st.session_state.insights = insights
                
                st.balloons()  # Celebrate success!
            else:
                st.error("❌ Could not generate charts. The dataset format may not be suitable for automatic analysis.")
        
        # Display charts
        if 'charts' in st.session_state and st.session_state.charts:
            st.header("📊 Smart Business Charts")
            
            # Create tabs for better organization
            chart_names = list(st.session_state.charts.keys())
            if len(chart_names) > 1:
                tabs = st.tabs([st.session_state.charts[name]['title'] for name in chart_names])
                
                for i, (chart_name, tab) in enumerate(zip(chart_names, tabs)):
                    with tab:
                        chart_info = st.session_state.charts[chart_name]
                        st.plotly_chart(
                            chart_info['figure'], 
                            use_container_width=True
                        )
                        
                        # Show corresponding insight
                        if 'insights' in st.session_state:
                            insight = next((ins for ins in st.session_state.insights 
                                          if ins['chart'] == chart_name), None)
                            if insight:
                                st.info(f"💡 **Insight:** {insight['insight']}")
            else:
                # Single chart
                for chart_name, chart_info in st.session_state.charts.items():
                    st.subheader(chart_info['title'])
                    st.plotly_chart(
                        chart_info['figure'], 
                        use_container_width=True
                    )
        
        # Generate PowerPoint
        if 'insights' in st.session_state and st.session_state.insights:
            st.header("📝 AI-Powered Insights")
            
            for insight in st.session_state.insights:
                with st.expander(f"💡 {insight['title']}", expanded=True):
                    st.write(insight['insight'])
            
            col1, col2 = st.columns([1, 1])
            with col1:
                if st.button("🎯 Generate PowerPoint Report", type="primary"):
                    with st.spinner("Creating PowerPoint with charts..."):
                        try:
                            filename = st.session_state.dashboard.create_powerpoint_with_charts(
                                st.session_state.charts, 
                                st.session_state.insights
                            )
                            st.success(f"✅ PowerPoint created: {filename}")
                            
                            # Provide download link
                            with open(filename, "rb") as file:
                                st.download_button(
                                    label="📥 Download PowerPoint",
                                    data=file.read(),
                                    file_name=filename.split('/')[-1],
                                    mime="application/vnd.ms-powerpoint"
                                )
                        except Exception as e:
                            st.error(f"Error creating PowerPoint: {str(e)}")
        
        # Chat with Data Section  
        st.header("💬 Chat with Your Data")
        
        if "chat_history" not in st.session_state:
            st.session_state.chat_history = []
        
        # Chat input
        user_question = st.chat_input("Ask anything about your data...")
        
        if user_question:
            # Add user message to chat history
            st.session_state.chat_history.append({"role": "user", "content": user_question})
            
            # Show what's happening
            status_placeholder = st.empty()
            
            # Get AI response
            with st.spinner("🤖 Analyzing your question..."):
                try:
                    # Get current ollama settings from sidebar
                    current_ollama_url = st.session_state.get('ollama_url', 'http://localhost:11434')
                    current_ollama_model = st.session_state.get('ollama_model', 'llama3.1')
                    
                    # Show status
                    if any(word in user_question.lower() for word in ['trend', 'pattern', 'insight', 'analysis']):
                        status_placeholder.info("🔍 Performing enhanced trend analysis...")
                    else:
                        status_placeholder.info("🤖 Getting instant response...")
                    
                    ai_response = chat_with_data(data, user_question, current_ollama_model, current_ollama_url)
                    st.session_state.chat_history.append({"role": "assistant", "content": ai_response})
                    
                    # Clear status
                    status_placeholder.empty()
                    
                except Exception as e:
                    status_placeholder.empty()
                    error_response = f"❌ Error processing your question: {str(e)}\n\nTry asking simpler questions like:\n• 'What columns are in this data?'\n• 'What are the key statistics?'\n• 'Tell me about this dataset'"
                    st.session_state.chat_history.append({"role": "assistant", "content": error_response})
        
        # Display chat history
        if st.session_state.chat_history:
            for message in st.session_state.chat_history[-8:]:  # Show last 8 messages
                with st.chat_message(message["role"]):
                    st.write(message["content"])
        else:
            # Show example questions and quick buttons if no chat history
            st.info("""
            💡 **Try asking questions like:**
            • "What are the key trends in this data?"
            • "Show me marketing performance insights"
            • "What columns have the highest values?"
            • "Tell me about data quality"
            • "What's the average CTR?"
            """)
            
            # Quick action buttons
            st.markdown("**🚀 Quick Questions:**")
            col1, col2, col3 = st.columns(3)
            
            with col1:
                if st.button("📊 Key Statistics", help="Get statistical summary"):
                    st.session_state.chat_history.append({"role": "user", "content": "What are the key statistics?"})
                    st.rerun()
            
            with col2:
                if st.button("📈 Trends", help="Analyze data trends"):
                    st.session_state.chat_history.append({"role": "user", "content": "What are the key trends in this data?"})
                    st.rerun()
                    
            with col3:
                if st.button("📋 Columns", help="Show all columns"):
                    st.session_state.chat_history.append({"role": "user", "content": "What columns are available?"})
                    st.rerun()
        
        # Clear chat button
        col1, col2 = st.columns([1, 4])
        with col1:
            if st.button("🗑️ Clear Chat"):
                st.session_state.chat_history = []
                st.rerun()
    
    else:
        # Welcome message
        st.info("👆 Load a dataset from the sidebar to get started!")
        
        st.markdown("""
        ### 🎯 Enhanced Features:
        
        1. **📤 CSV Upload**: Upload your own CSV files directly
        2. **📊 Smart Charts**: AI-generated business-relevant visualizations
        3. **🤖 Advanced Insights**: LLM-powered analysis and commentary
        4. **📑 Rich PowerPoint**: Auto-generate slides with embedded charts
        5. **💬 Data Chat**: Ask questions about your data in natural language
        
        ### 🚀 Getting Started:
        1. **Upload a CSV** or enter a **Kaggle dataset name**
        2. Set up **Ollama** for enhanced AI features (optional but recommended)
        3. Generate smart charts and insights
        4. Chat with your data using natural language
        5. Download comprehensive PowerPoint reports!
        
        ### 💡 Example Questions to Ask:
        - "What are the key trends in this data?"
        - "Which columns have the most missing values?"
        - "Show me the correlation between sales and profit"
        - "What insights can you derive from the patterns?"
        - "Which region or category needs the most attention?"
        
        ### 📊 Enhanced Chart Features:
        - **⚡ 10-15 Second Generation**: Up to 8 comprehensive charts with full dataset support  
        - **📍 High/Low Highlights**: Peak (📈) and trough (📉) points automatically marked
        - **🏆 Performance Indicators**: Best (🏆) and worst (⚠️) categories highlighted  
        - **🔧 Multiple Chart Types**: Time series, distributions, categories, correlations, geographic
        - **🎯 Smart Annotations**: Exact values displayed on all visualizations
        - **💬 Working Chatbot**: Ask questions and get instant answers about your data
        """)

if __name__ == "__main__":
    main()